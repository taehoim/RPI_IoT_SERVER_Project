"""
CM4 IoT Gateway 자체 호스팅 플랫폼 - 상세 구현 계획서 PPTX 생성기
- 38슬라이드, 16:9
- 시스템 개념도/구성도는 네이티브 shape로 작성 (PowerPoint에서 편집 가능)
- 폰트: Malgun Gothic (Windows 호환), 본문 한글 가독성 우선
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ===== Design tokens =====
NAVY        = RGBColor(0x1C, 0x2E, 0x4A)
DARK        = RGBColor(0x0F, 0x1A, 0x2A)
SLATE       = RGBColor(0x2E, 0x40, 0x53)
STEEL       = RGBColor(0x34, 0x49, 0x5E)
AMBER       = RGBColor(0xE6, 0x7E, 0x22)
GOLD        = RGBColor(0xF3, 0x9C, 0x12)
TEAL        = RGBColor(0x16, 0xA0, 0x85)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED         = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_BG    = RGBColor(0xF4, 0xF6, 0xF8)
SURFACE     = RGBColor(0xFF, 0xFF, 0xFF)
LINE        = RGBColor(0xBD, 0xC3, 0xC7)
MUTED       = RGBColor(0x7B, 0x8A, 0x8B)
TEXT_DARK   = RGBColor(0x1C, 0x28, 0x33)
TEXT_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE      = RGBColor(0xEC, 0xF0, 0xF1)

FONT_KO = "Malgun Gothic"
FONT_EN = "Arial"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ===== Helpers =====

def add_rect(slide, x, y, w, h, fill, line_color=None, line_w=0.5, shadow=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    if not shadow:
        # remove default shadow
        spPr = sh.fill._xPr
        sppr = sh._element.spPr
        for ef in sppr.findall(qn('a:effectLst')):
            sppr.remove(ef)
        eff = etree.SubElement(sppr, qn('a:effectLst'))
    sh.shadow.inherit = False
    return sh

def add_round_rect(slide, x, y, w, h, fill, line_color=None, line_w=0.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def set_text(shape, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             font=FONT_KO, anchor=MSO_ANCHOR.MIDDLE, italic=False):
    tf = shape.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, font=FONT_KO, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    set_text(tb, text, size=size, bold=bold, color=color, align=align, font=font, anchor=anchor, italic=italic)
    return tb

def add_multiline(slide, x, y, w, h, lines, size=12, color=TEXT_DARK, font=FONT_KO,
                  bullet=False, line_spacing=1.25):
    """
    lines: list of (text, opts_dict) or strings.
    opts_dict supports: size, bold, color, indent, italic, mono, align
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top  = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for entry in lines:
        if isinstance(entry, tuple):
            text, opts = entry
        else:
            text, opts = entry, {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        p.space_after = Pt(opts.get("space_after", 2))
        if "indent" in opts and opts["indent"]:
            p.level = opts["indent"]
        run = p.add_run()
        prefix = "• " if (bullet and not opts.get("nobullet")) else opts.get("prefix", "")
        run.text = prefix + text
        run.font.name = FONT_MONO if opts.get("mono") else font
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", color)
    return tb

def add_connector(slide, x1, y1, x2, y2, color=STEEL, width=1.25, arrow=True, dashed=False):
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        # add arrow end
        ln = line.line._get_or_add_ln()
        tail = ln.find(qn('a:tailEnd'))
        if tail is None:
            tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')
    return line

# ===== Slide layout helpers =====

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank

def add_header(slide, title, subtitle=None, page_no=None, total=None):
    # Top accent strip
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), NAVY)
    add_rect(slide, Inches(0), Inches(0.85), SLIDE_W, Inches(0.05), AMBER)
    add_textbox(slide, Inches(0.45), Inches(0.12), Inches(11), Inches(0.45),
                title, size=22, bold=True, color=TEXT_LIGHT)
    if subtitle:
        add_textbox(slide, Inches(0.45), Inches(0.48), Inches(11), Inches(0.32),
                    subtitle, size=11, color=RGBColor(0xCC, 0xD4, 0xDB))
    if page_no is not None and total is not None:
        add_textbox(slide, Inches(12.3), Inches(0.32), Inches(0.9), Inches(0.4),
                    f"{page_no:02d} / {total:02d}", size=11, color=RGBColor(0xCC, 0xD4, 0xDB),
                    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # Footer
    add_rect(slide, Inches(0), Inches(7.30), SLIDE_W, Inches(0.20), DARK)
    add_textbox(slide, Inches(0.45), Inches(7.32), Inches(11), Inches(0.16),
                "CM4 IoT Gateway 자체 호스팅 플랫폼 — 상세 구현 계획서 (2026-05-02)",
                size=8, color=RGBColor(0x99, 0xA3, 0xA8))

# ===== Build presentation =====

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

TOTAL = 38

# ---------- Slide 1: Title ----------
def slide_title():
    s = blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # Diagonal accent
    accent = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.5), Inches(0), Inches(3.83), Inches(7.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = DARK
    accent.line.fill.background()
    accent.shadow.inherit = False
    # Amber thin bar
    add_rect(s, Inches(0.6), Inches(2.6), Inches(0.8), Inches(0.1), AMBER)
    add_textbox(s, Inches(0.6), Inches(2.8), Inches(11), Inches(0.6),
                "CM4 IoT Gateway 자체 호스팅 플랫폼", size=36, bold=True, color=TEXT_LIGHT)
    add_textbox(s, Inches(0.6), Inches(3.5), Inches(11), Inches(0.7),
                "Industrial CM4 기반 Self-hosted IoT Fleet Management Platform",
                size=18, color=RGBColor(0xCC, 0xD4, 0xDB))
    add_textbox(s, Inches(0.6), Inches(4.4), Inches(11), Inches(0.5),
                "상세 구현 계획서 (Detailed Implementation Plan)",
                size=20, bold=True, color=AMBER)
    # Meta box
    add_rect(s, Inches(0.6), Inches(5.5), Inches(8.0), Inches(1.4), DARK, line_color=AMBER, line_w=0.75)
    add_multiline(s, Inches(0.85), Inches(5.6), Inches(7.7), Inches(1.3), [
        ("작성일      2026-05-02", {"color": TEXT_LIGHT, "size": 12}),
        ("구축 방식   Docker 미사용 · systemd 기반 자체 호스팅", {"color": TEXT_LIGHT, "size": 12}),
        ("핵심 구성   VerneMQ + PostgreSQL + Keycloak + FastAPI + React (Vite)", {"color": TEXT_LIGHT, "size": 12}),
        ("관리 모델   사용자별 다중 Gateway · Sensor Profile 기반 동적 관리", {"color": TEXT_LIGHT, "size": 12}),
    ], font=FONT_KO)
    add_textbox(s, Inches(0.6), Inches(7.0), Inches(11), Inches(0.3),
                "v1.0  ·  Owner: Solo + Claude Code AI pair", size=10, color=MUTED)

# ---------- Slide 2: Agenda ----------
def slide_agenda():
    s = blank_slide(prs)
    add_header(s, "목차 (Agenda)", "전체 41-section 계획서를 38 슬라이드로 압축", 2, TOTAL)
    items = [
        ("Part 1.  배경 및 설계 원칙", [
            "프로젝트 개요 / 핵심 요구사항 / 사업·기술 전제",
            "핵심 설계 철학 4가지",
        ]),
        ("Part 2.  시스템 아키텍처", [
            "전체 시스템 개념도",
            "상세 시스템 구성도 (서버 + Gateway 내부)",
            "서버 구성요소 / systemd 서비스 / 디렉터리 구조",
        ]),
        ("Part 3.  데이터 모델 및 통신", [
            "권한 모델 / Keycloak / PostgreSQL ERD",
            "Sensor Profile · Channel 분리 / Telemetry / Config 버전 관리",
            "MQTT Topic & Payload / 원격 제어 흐름 / Alarm Rule",
        ]),
        ("Part 4.  Web Portal · Gateway Agent · 보안", [
            "Web Portal 화면 / Sensor Wizard / Dashboard / Bulk",
            "Gateway Agent 구현 / Safety MCU 설계 / API 요약",
            "보안 설계 / RLS / 운영 관리",
        ]),
        ("Part 5.  실행 계획", [
            "개발 로드맵 7단계",
            "Phase 1 상세 spec (1주 sprint, 코드 0줄)",
            "라이선스 검토 / 우선순위 / 결론",
        ]),
    ]
    y = Inches(1.10)
    for title, subs in items:
        add_rect(s, Inches(0.45), y, Inches(0.08), Inches(1.05), AMBER)
        add_textbox(s, Inches(0.7), y + Inches(0.02), Inches(11.5), Inches(0.35),
                    title, size=14, bold=True, color=NAVY)
        sub_text = "   ·   ".join(subs)
        add_textbox(s, Inches(0.7), y + Inches(0.40), Inches(11.5), Inches(0.65),
                    sub_text, size=10.5, color=SLATE)
        y += Inches(1.20)

# ---------- Slide 3: 프로젝트 개요 ----------
def slide_overview():
    s = blank_slide(prs)
    add_header(s, "프로젝트 개요", "Industrial CM4 IoT Gateway Fleet Management Platform", 3, TOTAL)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
                "산업용 Raspberry Pi CM4 기반 IoT Gateway 제품을 자체 호스팅 서버와 연동하여 상용화한다.",
                size=14, bold=True, color=NAVY)
    add_textbox(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.4),
                "여러 사용자가 각자 하나 이상의 Gateway를 보유하며, Gateway마다 서로 다른 센서·제어 장치를 연결한다.",
                size=12, color=SLATE)
    # 10 requirements as 2x5 cards
    add_textbox(s, Inches(0.5), Inches(2.15), Inches(12), Inches(0.35),
                "10가지 핵심 요구사항", size=14, bold=True, color=AMBER)
    reqs = [
        "사용자별 여러 Gateway 관리",
        "Gateway별 서로 다른 센서 구성 관리",
        "릴레이·밸브·펌프 등 제어 채널 관리",
        "센서 종류 추가 시 코드 수정 최소화",
        "Gateway 설정 중앙 관리 및 원격 반영",
        "사용자·고객사·현장·장비 단위 권한",
        "원격 제어 명령의 안전성 확보",
        "장비 상태·센서·알람·이력 통합 관리",
        "Docker 없이 systemd 기반 운영",
        "무료 OSS 기반 상용화 가능 구조",
    ]
    cols = 2
    rows = 5
    cw = Inches(6.1)
    ch = Inches(0.85)
    x0 = Inches(0.5)
    y0 = Inches(2.55)
    for i, r in enumerate(reqs):
        col = i % cols
        row = i // cols
        x = x0 + col * (cw + Inches(0.1))
        y = y0 + row * (ch + Inches(0.08))
        add_round_rect(s, x, y, cw, ch, SURFACE, line_color=LINE, line_w=0.5)
        # Number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), y + Inches(0.16), Inches(0.55), Inches(0.55))
        badge.fill.solid(); badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background(); badge.shadow.inherit = False
        set_text(badge, f"{i+1:02d}", size=12, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.8), y + Inches(0.17), cw - Inches(0.9), Inches(0.55),
                    r, size=12, bold=True, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

# ---------- Slide 4: 사업/기술 전제 ----------
def slide_premises():
    s = blank_slide(prs)
    add_header(s, "전제 조건", "사업·운영 전제 / 기술 스택 (인터뷰 결과 반영)", 4, TOTAL)
    # Two columns
    # Left: business
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(5.9), SURFACE, line_color=LINE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(0.5), NAVY)
    add_textbox(s, Inches(0.7), Inches(1.13), Inches(5.5), Inches(0.45),
                "사업 · 운영 전제", size=14, bold=True, color=TEXT_LIGHT)
    biz = [
        ("서비스 운영", "자체 서버 운영 (퍼블릭 IoT 플랫폼 미사용)"),
        ("컨테이너", "Docker 미사용"),
        ("서버 OS", "Ubuntu Server 24.04 LTS · 사내 물리 머신"),
        ("사용자 구조", "한 사용자가 여러 Gateway 보유 가능"),
        ("센서 구성", "Gateway마다 종류·수량 상이"),
        ("제어 구성", "Gateway마다 릴레이·밸브·펌프 구성 상이"),
        ("관리회사", "전체 사용자·Gateway 통합 관리"),
        ("일반 사용자", "본인에게 할당된 Gateway만 접근"),
        ("라이선스", "무료 OSS 중심 (AGPL/BSL/GPL 리스크 최소화)"),
    ]
    y = Inches(1.75)
    for k, v in biz:
        add_textbox(s, Inches(0.7), y, Inches(1.6), Inches(0.32), k, size=10.5, bold=True, color=AMBER)
        add_textbox(s, Inches(2.3), y, Inches(4.2), Inches(0.32), v, size=10.5, color=TEXT_DARK)
        y += Inches(0.55)
    # Right: tech
    add_round_rect(s, Inches(6.8), Inches(1.1), Inches(6.1), Inches(5.9), SURFACE, line_color=LINE)
    add_rect(s, Inches(6.8), Inches(1.1), Inches(6.1), Inches(0.5), STEEL)
    add_textbox(s, Inches(7.0), Inches(1.13), Inches(5.5), Inches(0.45),
                "기술 스택 (확정)", size=14, bold=True, color=TEXT_LIGHT)
    tech = [
        ("MQTT Broker", "VerneMQ (Apache 2.0)"),
        ("DB", "PostgreSQL 16 (apt)"),
        ("인증", "Keycloak — Postgres backend"),
        ("Backend", "FastAPI (Python 3.12)"),
        ("Frontend", "React + Vite (SPA)"),
        ("Chart", "Apache ECharts"),
        ("Reverse Proxy", "Nginx + certbot"),
        ("서비스 관리", "systemd"),
        ("Gateway OS", "Pi OS Lite / Ubuntu ARM64 (Phase 2 결정)"),
        ("Gateway 통신", "MQTT (Phase 1: plain · Phase 7: TLS)"),
        ("Gateway 설정", "서버 desired/reported 구조"),
    ]
    y = Inches(1.75)
    for k, v in tech:
        add_textbox(s, Inches(7.0), y, Inches(1.7), Inches(0.30), k, size=10.5, bold=True, color=AMBER)
        add_textbox(s, Inches(8.7), y, Inches(4.1), Inches(0.30), v, size=10.5, color=TEXT_DARK)
        y += Inches(0.48)

# ---------- Slide 5: 핵심 설계 철학 ----------
def slide_philosophy():
    s = blank_slide(prs)
    add_header(s, "핵심 설계 철학", "이 4가지 원칙이 모든 데이터 모델·코드 구조를 결정한다", 5, TOTAL)
    items = [
        ("01", "Gateway 중심 엔티티",
         "사용자별 여러 Gateway, Gateway별 회사·현장 배정. User-Company-Site-Gateway 4축으로 권한 판단."),
        ("02", "Sensor 종류는 코드가 아닌 Profile",
         "센서 추가 시 schema/Gateway 코드 수정 없음. Sensor Profile (모델 정의) + Sensor Channel (실 인스턴스) 분리."),
        ("03", "Gateway Config 중앙 버전 관리",
         "서버가 Gateway별 설정을 desired_config로 발행 → Gateway가 reported_config로 보고. config_version + config_hash 추적."),
        ("04", "systemd 기반 서비스 운영",
         "Docker 미사용. nginx · vernemq · postgresql · keycloak · iot-backend · iot-worker · iot-scheduler를 OS 서비스로 명확 관리."),
    ]
    cw = Inches(6.05); ch = Inches(2.9)
    coords = [
        (Inches(0.5), Inches(1.1)),
        (Inches(6.78), Inches(1.1)),
        (Inches(0.5), Inches(4.15)),
        (Inches(6.78), Inches(4.15)),
    ]
    for (no, title, desc), (x, y) in zip(items, coords):
        add_round_rect(s, x, y, cw, ch, SURFACE, line_color=LINE)
        add_rect(s, x, y, Inches(0.15), ch, AMBER)
        add_textbox(s, x + Inches(0.4), y + Inches(0.15), Inches(1.0), Inches(0.5),
                    no, size=28, bold=True, color=AMBER)
        add_textbox(s, x + Inches(1.4), y + Inches(0.25), cw - Inches(1.6), Inches(0.5),
                    title, size=15, bold=True, color=NAVY)
        add_textbox(s, x + Inches(0.4), y + Inches(1.0), cw - Inches(0.7), ch - Inches(1.1),
                    desc, size=12, color=TEXT_DARK)

# ---------- Slide 6: 전체 시스템 개념도 ----------
def slide_concept_diagram():
    s = blank_slide(prs)
    add_header(s, "전체 시스템 개념도", "Concept Architecture — 사용자, 서버, Gateway, 현장 장비의 관계", 6, TOTAL)

    # Layer bands
    # User layer
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.0), SUBTLE, line_color=LINE)
    add_textbox(s, Inches(0.6), Inches(1.15), Inches(2.0), Inches(0.3), "USER LAYER",
                size=10, bold=True, color=MUTED)
    # 3 user types
    user_types = [("일반 사용자", "본인 Gateway 조회·제어"), ("고객사 관리자", "회사 내 사용자·현장 관리"), ("관리회사 관리자", "전체 통합 관리")]
    for i, (n, d) in enumerate(user_types):
        x = Inches(1.0 + i * 4.0)
        b = add_round_rect(s, x, Inches(1.45), Inches(3.5), Inches(0.6), NAVY)
        set_text(b, n, size=12, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(2.05), Inches(3.5), Inches(0.2), d,
                    size=8.5, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Web Portal layer
    web = add_round_rect(s, Inches(4.5), Inches(2.5), Inches(4.3), Inches(0.65), STEEL)
    set_text(web, "Web Portal  (React + Vite SPA · ECharts)", size=12, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Server box (large)
    add_round_rect(s, Inches(0.5), Inches(3.4), Inches(8.5), Inches(2.7), SURFACE, line_color=NAVY, line_w=1.5)
    add_textbox(s, Inches(0.7), Inches(3.45), Inches(7), Inches(0.3), "SELF-HOSTED SERVER  (Ubuntu 24.04 LTS · systemd)",
                size=10, bold=True, color=NAVY)
    # Inner server services
    services = [
        ("Nginx", "443/80", AMBER),
        ("Keycloak", "OIDC", TEAL),
        ("Backend API", "FastAPI", NAVY),
        ("Worker", "MQTT 처리", SLATE),
        ("Scheduler", "주기 작업", SLATE),
        ("PostgreSQL", "16", GREEN),
        ("VerneMQ", "1883/8883", AMBER),
    ]
    sx = Inches(0.7); sy = Inches(3.85)
    for i, (name, sub, col) in enumerate(services):
        col_n = i % 4
        row_n = i // 4
        x = sx + col_n * Inches(2.05)
        y = sy + row_n * Inches(1.05)
        b = add_round_rect(s, x, y, Inches(1.95), Inches(0.9), col)
        add_textbox(s, x, y + Inches(0.08), Inches(1.95), Inches(0.4), name,
                    size=12, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(s, x, y + Inches(0.45), Inches(1.95), Inches(0.4), sub,
                    size=10, color=RGBColor(0xEC, 0xF0, 0xF1), align=PP_ALIGN.CENTER)

    # MQTT Broker (separate emphasis on right)
    # File storage
    add_round_rect(s, Inches(9.2), Inches(3.4), Inches(3.6), Inches(2.7), DARK, line_color=AMBER, line_w=1)
    add_textbox(s, Inches(9.35), Inches(3.45), Inches(3.4), Inches(0.3), "STORAGE / ASSETS",
                size=10, bold=True, color=AMBER)
    storage_items = [("펌웨어 저장소", "/var/lib/.../firmware"),
                     ("Gateway Configs", "/var/lib/.../gateway-configs"),
                     ("백업", "/var/lib/.../backups"),
                     ("로그 번들", "/var/lib/.../log-bundles")]
    sy2 = Inches(3.85)
    for n, p in storage_items:
        add_textbox(s, Inches(9.35), sy2, Inches(3.3), Inches(0.25), n,
                    size=11, bold=True, color=TEXT_LIGHT)
        add_textbox(s, Inches(9.35), sy2 + Inches(0.25), Inches(3.3), Inches(0.22), p,
                    size=9, color=RGBColor(0xCC, 0xD4, 0xDB), font=FONT_MONO)
        sy2 += Inches(0.55)

    # Gateway layer
    add_rect(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.85), SUBTLE, line_color=LINE)
    add_textbox(s, Inches(0.6), Inches(6.34), Inches(3.0), Inches(0.3),
                "GATEWAY LAYER  (CM4 + Safety MCU)", size=10, bold=True, color=MUTED)
    # Multiple gateways
    for i in range(4):
        x = Inches(1.5 + i * 2.85)
        gb = add_round_rect(s, x, Inches(6.62), Inches(2.55), Inches(0.45), AMBER)
        set_text(gb, f"CM4 Gateway #{i+1:02d}", size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Connectors
    # Users -> Web Portal (3 lines)
    add_connector(s, Inches(2.75), Inches(2.10), Inches(6.65), Inches(2.50), color=MUTED, width=1)
    add_connector(s, Inches(6.75), Inches(2.10), Inches(6.65), Inches(2.50), color=MUTED, width=1)
    add_connector(s, Inches(10.75), Inches(2.10), Inches(6.65), Inches(2.50), color=MUTED, width=1)
    # Web Portal -> Server (HTTPS)
    add_connector(s, Inches(6.65), Inches(3.15), Inches(4.75), Inches(3.40), color=NAVY, width=1.5)
    # Server -> Gateways (MQTT TLS)
    for i in range(4):
        x = Inches(2.5 + i * 2.85) + Inches(0.3)
        add_connector(s, Inches(8.0), Inches(6.10), x, Inches(6.62), color=AMBER, width=1.25, dashed=True)

    # Legend
    add_textbox(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.2),
                "── HTTPS (Web/API)     ┄┄ MQTT (Telemetry · Command · Config)",
                size=9, color=MUTED, font=FONT_MONO)

# ---------- Slide 7: 상세 시스템 구성도 ----------
def slide_detailed_arch():
    s = blank_slide(prs)
    add_header(s, "상세 시스템 구성도", "Detailed Architecture — 서버 내부 모듈, 포트, 데이터 흐름, Gateway 내부 구조", 7, TOTAL)

    # ---- Server side (left) ----
    add_round_rect(s, Inches(0.4), Inches(1.05), Inches(7.6), Inches(6.1), SURFACE, line_color=NAVY, line_w=1.5)
    add_rect(s, Inches(0.4), Inches(1.05), Inches(7.6), Inches(0.4), NAVY)
    add_textbox(s, Inches(0.55), Inches(1.07), Inches(7.0), Inches(0.36),
                "서버 (Ubuntu 24.04 · systemd)", size=12, bold=True, color=TEXT_LIGHT)

    # Nginx
    n = add_round_rect(s, Inches(0.7), Inches(1.65), Inches(7.0), Inches(0.55), AMBER)
    set_text(n, "Nginx  ·  443/80  ·  HTTPS reverse proxy + 정적 파일",
             size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # Routes under nginx
    add_textbox(s, Inches(0.7), Inches(2.25), Inches(7.0), Inches(0.25),
                "/ → frontend     /api/ → 8000     /auth/ → 8080",
                size=9, color=MUTED, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # 3 main app services row
    apps = [("iot-backend.service", "FastAPI · :8000\n• REST API\n• 권한 검사\n• Command publish\n• Config 생성"),
            ("iot-worker.service", "Python · MQTT\n• telemetry 수신\n• telemetry_latest\n• alarm eval\n• response 처리"),
            ("iot-scheduler.service", "Python · 주기\n• offline 판정\n• cmd timeout\n• partition 정리\n• report 생성")]
    for i, (name, body) in enumerate(apps):
        x = Inches(0.7 + i * 2.42)
        add_round_rect(s, x, Inches(2.6), Inches(2.32), Inches(1.55), STEEL)
        add_textbox(s, x, Inches(2.65), Inches(2.32), Inches(0.32), name,
                    size=10, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font=FONT_MONO)
        add_textbox(s, x + Inches(0.1), Inches(2.97), Inches(2.12), Inches(1.18), body,
                    size=8.5, color=RGBColor(0xEC, 0xF0, 0xF1))

    # Auth + DB + MQTT row (infra)
    infra = [
        ("Keycloak", ":8080", "OIDC/OAuth2 · realm: iot-platform · roles 7", TEAL),
        ("PostgreSQL", ":5432", "iot_platform · keycloak  ·  RLS 검토", GREEN),
        ("VerneMQ", ":1883/:8883", "MQTT broker  ·  password→X.509 (Phase 7)", AMBER),
    ]
    for i, (n, port, desc, col) in enumerate(infra):
        x = Inches(0.7 + i * 2.42)
        add_round_rect(s, x, Inches(4.3), Inches(2.32), Inches(1.4), col)
        add_textbox(s, x, Inches(4.35), Inches(2.32), Inches(0.32), n,
                    size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(4.65), Inches(2.32), Inches(0.25), port,
                    size=9, color=RGBColor(0xEC, 0xF0, 0xF1), align=PP_ALIGN.CENTER, font=FONT_MONO)
        add_textbox(s, x + Inches(0.1), Inches(4.95), Inches(2.12), Inches(0.7), desc,
                    size=8.5, color=RGBColor(0xEC, 0xF0, 0xF1), align=PP_ALIGN.CENTER)

    # Storage row
    add_round_rect(s, Inches(0.7), Inches(5.85), Inches(7.0), Inches(1.2), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.85), Inches(5.9), Inches(6.7), Inches(0.3),
                "/var/lib/iot-platform/  ·  /etc/iot-platform/  ·  /var/log/iot-platform/",
                size=10, bold=True, color=AMBER, font=FONT_MONO)
    storage = [("firmware/", "OTA 이미지"), ("gateway-configs/", "config 백업"),
               ("backups/", "pg_dump"), ("log-bundles/", "장비 로그 수집")]
    for i, (k, v) in enumerate(storage):
        x = Inches(0.85 + (i % 4) * 1.7)
        add_textbox(s, x, Inches(6.2), Inches(1.6), Inches(0.25), k,
                    size=9, bold=True, color=TEXT_LIGHT, font=FONT_MONO)
        add_textbox(s, x, Inches(6.45), Inches(1.6), Inches(0.5), v,
                    size=8.5, color=RGBColor(0xCC, 0xD4, 0xDB))

    # ---- Gateway side (right) ----
    add_round_rect(s, Inches(8.2), Inches(1.05), Inches(4.8), Inches(6.1), SURFACE, line_color=AMBER, line_w=1.5)
    add_rect(s, Inches(8.2), Inches(1.05), Inches(4.8), Inches(0.4), AMBER)
    add_textbox(s, Inches(8.35), Inches(1.07), Inches(4.5), Inches(0.36),
                "Gateway (CM4 Linux + Safety MCU)", size=12, bold=True, color=DARK)

    # CM4 Linux modules
    add_round_rect(s, Inches(8.4), Inches(1.6), Inches(4.4), Inches(3.4), STEEL)
    add_textbox(s, Inches(8.5), Inches(1.65), Inches(4.2), Inches(0.3),
                "/opt/iot-gateway/   (Linux user space)",
                size=10, bold=True, color=TEXT_LIGHT, font=FONT_MONO)
    gw_mods = [
        ("gateway-agent", "MQTT 연결 · 부팅 흐름"),
        ("sensor-service", "Modbus/AI/DI polling"),
        ("actuator-service", "GPIO/Relay 제어"),
        ("rule-engine", "Local Rule (오프라인 가용)"),
        ("mqtt-client", "Telemetry/Command/Config"),
        ("local-db (SQLite)", "telemetry queue · cmd log"),
        ("ota-agent", "이미지 검증·적용"),
        ("health-agent", "CPU/MEM/Net heartbeat"),
    ]
    for i, (n, d) in enumerate(gw_mods):
        col = i % 2
        row = i // 2
        x = Inches(8.5 + col * 2.15)
        y = Inches(2.0 + row * 0.7)
        add_round_rect(s, x, y, Inches(2.05), Inches(0.6), DARK)
        add_textbox(s, x + Inches(0.05), y + Inches(0.04), Inches(2), Inches(0.28),
                    n, size=9, bold=True, color=AMBER, font=FONT_MONO)
        add_textbox(s, x + Inches(0.05), y + Inches(0.32), Inches(2), Inches(0.28),
                    d, size=8.5, color=TEXT_LIGHT)

    # Safety MCU
    add_round_rect(s, Inches(8.4), Inches(5.1), Inches(4.4), Inches(0.85), RED)
    add_textbox(s, Inches(8.5), Inches(5.15), Inches(4.2), Inches(0.3),
                "STM32 / NXP Safety MCU", size=11, bold=True, color=TEXT_LIGHT)
    add_textbox(s, Inches(8.5), Inches(5.45), Inches(4.2), Inches(0.45),
                "릴레이 직접 제어 · Watchdog · Fail-safe · Interlock · 비상 정지",
                size=9, color=TEXT_LIGHT)

    # Field IO
    add_round_rect(s, Inches(8.4), Inches(6.05), Inches(4.4), Inches(1.0), SUBTLE, line_color=LINE)
    add_textbox(s, Inches(8.5), Inches(6.1), Inches(4.2), Inches(0.3),
                "Field I/O", size=10, bold=True, color=NAVY)
    add_textbox(s, Inches(8.5), Inches(6.4), Inches(4.2), Inches(0.6),
                "RS-485 (Modbus RTU)  ·  GPIO Relay  ·  Analog 0-10V/4-20mA  ·  GPS  ·  Digital Input",
                size=9, color=TEXT_DARK)

    # Cross connector — server VerneMQ to Gateway MQTT client
    add_connector(s, Inches(7.7), Inches(4.85), Inches(8.5), Inches(2.5),
                  color=AMBER, width=2, dashed=True)
    add_textbox(s, Inches(7.65), Inches(3.5), Inches(0.55), Inches(0.4),
                "MQTT", size=9, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ---------- Slide 8: systemd 서비스 구성 ----------
def slide_systemd():
    s = blank_slide(prs)
    add_header(s, "서버 구성 — systemd 서비스 6+종", "Docker 미사용. OS 서비스 단위로 명확 분리", 8, TOTAL)
    rows = [
        ("nginx.service", "필수", "HTTPS reverse proxy · 정적 파일 제공"),
        ("vernemq.service", "필수", "MQTT Broker (1883 plain · 추후 8883 TLS)"),
        ("postgresql.service", "필수", "관계형 DB · 센서 데이터 저장 · keycloak DB 공유"),
        ("keycloak.service", "필수", "사용자 인증 · OIDC/OAuth2 토큰 발급"),
        ("iot-backend.service", "필수", "REST API · 권한 검사 · Gateway 관리 · 명령 발행"),
        ("iot-worker.service", "필수", "MQTT subscribe · telemetry ingestion · alarm eval"),
        ("iot-scheduler.service", "필수", "offline 판단 · timeout 처리 · 백업 · OTA 상태"),
        ("prometheus.service", "선택", "내부 metric 수집 (Phase 7+)"),
        ("opensearch.service", "선택", "로그 검색 (필요 시)"),
    ]
    # Header
    x0 = Inches(0.5); y0 = Inches(1.15); w = Inches(12.3)
    add_rect(s, x0, y0, w, Inches(0.45), NAVY)
    add_textbox(s, x0 + Inches(0.15), y0 + Inches(0.06), Inches(3), Inches(0.32),
                "서비스 이름", size=11, bold=True, color=TEXT_LIGHT)
    add_textbox(s, x0 + Inches(3.5), y0 + Inches(0.06), Inches(1), Inches(0.32),
                "분류", size=11, bold=True, color=TEXT_LIGHT)
    add_textbox(s, x0 + Inches(4.7), y0 + Inches(0.06), Inches(7.5), Inches(0.32),
                "역할", size=11, bold=True, color=TEXT_LIGHT)
    y = y0 + Inches(0.45)
    for name, kind, desc in rows:
        bg = SURFACE if (rows.index((name, kind, desc)) % 2 == 0) else SUBTLE
        add_rect(s, x0, y, w, Inches(0.5), bg, line_color=LINE, line_w=0.25)
        add_textbox(s, x0 + Inches(0.15), y + Inches(0.08), Inches(3.3), Inches(0.34),
                    name, size=11, bold=True, color=NAVY, font=FONT_MONO)
        col = AMBER if kind == "필수" else MUTED
        add_textbox(s, x0 + Inches(3.5), y + Inches(0.08), Inches(1), Inches(0.34),
                    kind, size=10, bold=True, color=col)
        add_textbox(s, x0 + Inches(4.7), y + Inches(0.08), Inches(7.5), Inches(0.34),
                    desc, size=10.5, color=TEXT_DARK)
        y += Inches(0.5)
    # Note
    add_textbox(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                "※ 모든 서비스는 EnvironmentFile=/etc/iot-platform/{name}.env 로 환경 변수 분리. Restart=always · After=network.target.",
                size=9.5, italic=True, color=MUTED)

# ---------- Slide 9: 디렉터리 구조 ----------
def slide_directory():
    s = blank_slide(prs)
    add_header(s, "권장 서버 디렉터리 구조", "FHS 친화 · 환경 변수 분리 · 릴리스 슬롯 + symlink", 9, TOTAL)
    blocks = [
        ("/opt/iot-platform/", NAVY, [
            "backend/   app/  venv/  migrations/  scripts/",
            "worker/    app/  venv/",
            "scheduler/ app/  venv/",
            "frontend/  current/  releases/",
            "releases/  backend-1.0.0/  backend-1.0.1/  ...",
            "current -> releases/backend-1.0.x   (atomic switch)",
        ]),
        ("/etc/iot-platform/", AMBER, [
            "backend.env       (DATABASE_URL, KC_ISSUER, ...)",
            "worker.env        (MQTT_HOST, MQTT_USER, ...)",
            "scheduler.env",
            "mqtt.env / db.env / keycloak.env",
        ]),
        ("/var/lib/iot-platform/", TEAL, [
            "firmware/         (OTA 이미지)",
            "gateway-configs/  (desired/reported snapshots)",
            "reports/          (정기 통계)",
            "log-bundles/      (Gateway 진단 로그)",
            "backups/          (pg_dump · realm export)",
        ]),
        ("/var/log/iot-platform/", STEEL, [
            "backend.log / worker.log / scheduler.log",
            "mqtt-ingestion.log",
        ]),
    ]
    cw = Inches(6.1); ch = Inches(2.85)
    coords = [
        (Inches(0.5), Inches(1.1)),
        (Inches(6.78), Inches(1.1)),
        (Inches(0.5), Inches(4.1)),
        (Inches(6.78), Inches(4.1)),
    ]
    for (label, color, lines), (x, y) in zip(blocks, coords):
        add_round_rect(s, x, y, cw, ch, SURFACE, line_color=LINE)
        add_rect(s, x, y, cw, Inches(0.4), color)
        add_textbox(s, x + Inches(0.15), y + Inches(0.05), cw - Inches(0.2), Inches(0.32),
                    label, size=12, bold=True, color=TEXT_LIGHT, font=FONT_MONO)
        items = [(l, {"mono": True, "size": 10, "color": TEXT_DARK}) for l in lines]
        add_multiline(s, x + Inches(0.2), y + Inches(0.5), cw - Inches(0.4), ch - Inches(0.55),
                      items, line_spacing=1.3)

# ---------- Slide 10: 역할 분담 ----------
def slide_roles():
    s = blank_slide(prs)
    add_header(s, "Backend / Worker / Scheduler 역할 분담", "관심사 분리 — 각 프로세스가 다른 라이프사이클을 가짐", 10, TOTAL)
    cards = [
        ("Backend API", NAVY, "REST · 동기 요청 처리", [
            "사용자 인증 연동 (Keycloak JWT 검증)",
            "User · Company · Site · Gateway 권한 검사",
            "Gateway 등록 / Profile / Channel 관리",
            "Gateway Config 생성 · 버전 관리",
            "MQTT command publish (request 발행)",
            "Telemetry 조회 API · latest 조회",
            "Alarm Rule 관리 · Audit Log",
        ]),
        ("Worker", STEEL, "MQTT subscribe · 비동기", [
            "MQTT subscribe (gw/+/telemetry, state, ...)",
            "Telemetry message 검증 (schema · ts)",
            "Telemetry 저장 (partition table)",
            "telemetry_latest upsert",
            "Heartbeat / state 처리",
            "Command response 처리",
            "Reported_config 처리 · Alarm evaluation",
        ]),
        ("Scheduler", AMBER, "주기 작업 · cron 대체", [
            "Gateway offline 판단 (heartbeat timeout)",
            "미응답 command timeout 처리",
            "오래된 telemetry partition drop",
            "백업 실행 (pg_dump · realm export)",
            "OTA job 상태 확인 · 재시도",
            "알람 재전송",
            "통계 / 리포트 생성",
        ]),
    ]
    cw = Inches(4.15); ch = Inches(5.95)
    for i, (name, color, desc, items) in enumerate(cards):
        x = Inches(0.5 + i * 4.27)
        y = Inches(1.1)
        add_round_rect(s, x, y, cw, ch, SURFACE, line_color=LINE)
        add_rect(s, x, y, cw, Inches(0.7), color)
        add_textbox(s, x + Inches(0.2), y + Inches(0.07), cw, Inches(0.36),
                    name, size=15, bold=True, color=TEXT_LIGHT)
        add_textbox(s, x + Inches(0.2), y + Inches(0.42), cw, Inches(0.26),
                    desc, size=10, color=RGBColor(0xEC, 0xF0, 0xF1))
        bullets = [(t, {"size": 10.5}) for t in items]
        add_multiline(s, x + Inches(0.25), y + Inches(0.85), cw - Inches(0.4), ch - Inches(0.95),
                      bullets, bullet=True, line_spacing=1.3)

# ---------- Slide 11: 권한 모델 계층 ----------
def slide_rbac_hierarchy():
    s = blank_slide(prs)
    add_header(s, "사용자 · 권한 모델", "Keycloak (인증) + Backend DB (실권한) — RBAC × ABAC", 11, TOTAL)
    # 좌측: 계층도
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.95), SURFACE, line_color=LINE)
    add_textbox(s, Inches(0.7), Inches(1.2), Inches(5.7), Inches(0.4),
                "권한 계층 (Top-Down)", size=13, bold=True, color=NAVY)
    levels = [
        ("System Admin",            "시스템 전체 설정 · 모든 회사·Gateway"),
        ("Management Company Admin","모든 고객사·전체 Gateway 관제·유지보수"),
        ("Customer Company Admin",  "본인 회사 사용자·현장·Gateway 관리"),
        ("Site Manager",            "특정 현장 Gateway 관리"),
        ("Operator",                "허용된 Gateway 제어 채널 조작"),
        ("Viewer",                  "센서 데이터·상태 조회만 가능"),
        ("Maintenance Engineer",    "진단·로그 수집·OTA·재부팅 (별도 축)"),
    ]
    y = Inches(1.65)
    for i, (name, desc) in enumerate(levels):
        # widening indent from top
        indent = Inches(0.0 + min(i, 5) * 0.3)
        is_aux = (name == "Maintenance Engineer")
        col = TEAL if is_aux else (NAVY if i == 0 else STEEL)
        add_round_rect(s, Inches(0.7) + indent, y, Inches(5.0) - indent, Inches(0.6), col)
        add_textbox(s, Inches(0.85) + indent, y + Inches(0.08), Inches(2.6), Inches(0.2),
                    name, size=11, bold=True, color=TEXT_LIGHT)
        add_textbox(s, Inches(0.85) + indent, y + Inches(0.32), Inches(4.7) - indent, Inches(0.25),
                    desc, size=9, color=RGBColor(0xEC, 0xF0, 0xF1))
        y += Inches(0.7)

    # 우측: 권한 테이블 + 액션 종류
    add_round_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(2.8), SURFACE, line_color=LINE)
    add_textbox(s, Inches(6.95), Inches(1.2), Inches(5.7), Inches(0.4),
                "권한 테이블 (3축)", size=13, bold=True, color=NAVY)
    perm_tables = [
        ("user_company_roles", "user × company × role"),
        ("user_site_permissions", "user × site × permission"),
        ("user_gateway_permissions", "user × gateway × permission"),
    ]
    yy = Inches(1.7)
    for n, d in perm_tables:
        add_textbox(s, Inches(6.95), yy, Inches(3.2), Inches(0.3), n,
                    size=11, bold=True, color=AMBER, font=FONT_MONO)
        add_textbox(s, Inches(10.2), yy, Inches(2.5), Inches(0.3), d,
                    size=10, color=TEXT_DARK)
        yy += Inches(0.55)
    add_textbox(s, Inches(6.95), Inches(3.4), Inches(5.7), Inches(0.4),
                "권한 종류 (permission)", size=11, bold=True, color=NAVY)
    add_textbox(s, Inches(6.95), Inches(3.65), Inches(5.7), Inches(0.3),
                "view  ·  control  ·  configure  ·  maintain  ·  admin",
                size=11, color=TEXT_DARK, font=FONT_MONO)

    # 우측 하단: Keycloak vs Backend
    add_round_rect(s, Inches(6.78), Inches(4.0), Inches(6.05), Inches(3.05), SURFACE, line_color=LINE)
    add_textbox(s, Inches(6.95), Inches(4.1), Inches(5.7), Inches(0.4),
                "Keycloak ↔ Backend 책임 분리", size=13, bold=True, color=NAVY)
    add_round_rect(s, Inches(6.95), Inches(4.55), Inches(2.8), Inches(2.4), TEAL)
    add_textbox(s, Inches(7.05), Inches(4.6), Inches(2.6), Inches(0.32),
                "Keycloak", size=12, bold=True, color=TEXT_LIGHT)
    add_textbox(s, Inches(7.05), Inches(4.95), Inches(2.6), Inches(1.95),
                "사용자 인증\nOIDC / OAuth2 토큰 발급\nrole claim 제공\n사용자 가입·비번 reset",
                size=10, color=TEXT_LIGHT)
    add_round_rect(s, Inches(9.95), Inches(4.55), Inches(2.85), Inches(2.4), NAVY)
    add_textbox(s, Inches(10.05), Inches(4.6), Inches(2.6), Inches(0.32),
                "Backend DB", size=12, bold=True, color=TEXT_LIGHT)
    add_textbox(s, Inches(10.05), Inches(4.95), Inches(2.7), Inches(1.95),
                "company_id · site_id · gateway_id 기반 실권한 판단\n사용자 ↔ 자산 매핑\n명령 권한 검증\nAudit Log",
                size=10, color=TEXT_LIGHT)

# ---------- Slide 12: Keycloak 구성 ----------
def slide_keycloak():
    s = blank_slide(prs)
    add_header(s, "Keycloak 구성", "Realm: iot-platform · 7 roles · group 기반 회사 구분", 12, TOTAL)
    # Realm card
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.7), TEAL)
    add_textbox(s, Inches(0.7), Inches(1.18), Inches(12), Inches(0.5),
                "Realm:  iot-platform        ·        DB Backend: PostgreSQL (공유 인스턴스, db=keycloak)",
                size=14, bold=True, color=TEXT_LIGHT)
    # Groups
    add_round_rect(s, Inches(0.5), Inches(1.95), Inches(6.05), Inches(2.55), SURFACE, line_color=LINE)
    add_textbox(s, Inches(0.7), Inches(2.05), Inches(5.5), Inches(0.4),
                "Groups (Company 매핑)", size=13, bold=True, color=NAVY)
    groups = ["management-company", "customer-company-a", "customer-company-b", "customer-company-c"]
    yy = Inches(2.5)
    for g in groups:
        add_round_rect(s, Inches(0.7), yy, Inches(5.7), Inches(0.4), SUBTLE, line_color=LINE)
        add_textbox(s, Inches(0.85), yy + Inches(0.05), Inches(5.4), Inches(0.32),
                    g, size=11, bold=True, color=NAVY, font=FONT_MONO)
        yy += Inches(0.48)

    # Roles
    add_round_rect(s, Inches(6.78), Inches(1.95), Inches(6.05), Inches(2.55), SURFACE, line_color=LINE)
    add_textbox(s, Inches(6.95), Inches(2.05), Inches(5.5), Inches(0.4),
                "Roles (7종, Phase 1에서 모두 생성)", size=13, bold=True, color=NAVY)
    roles = ["system_admin", "management_admin", "company_admin", "site_manager",
             "operator", "viewer", "maintenance_engineer"]
    cols = 2
    for i, r in enumerate(roles):
        c = i % cols; rr = i // cols
        x = Inches(6.95 + c * 2.85); y = Inches(2.5 + rr * 0.48)
        add_round_rect(s, x, y, Inches(2.7), Inches(0.4), AMBER)
        add_textbox(s, x + Inches(0.1), y + Inches(0.05), Inches(2.5), Inches(0.32),
                    r, size=10.5, bold=True, color=DARK, font=FONT_MONO)

    # Design rule
    add_round_rect(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.4), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.7), Inches(4.75), Inches(12), Inches(0.4),
                "설계 원칙", size=13, bold=True, color=AMBER)
    add_multiline(s, Inches(0.8), Inches(5.15), Inches(12), Inches(1.85), [
        ("Keycloak 책임", {"bold": True, "color": TEXT_LIGHT, "size": 11}),
        ("사용자 인증 · 비밀번호 정책 · OTP/2FA · 비밀번호 reset · 토큰 발급 · role claim 제공", {"color": RGBColor(0xCC, 0xD4, 0xDB), "size": 11}),
        ("Backend 책임", {"bold": True, "color": TEXT_LIGHT, "size": 11, "space_after": 4}),
        ("company_id · site_id · gateway_id 기반 실 접근 권한 판단 (Keycloak claim만 신뢰하지 않음)", {"color": RGBColor(0xCC, 0xD4, 0xDB), "size": 11}),
        ("토큰 검증 흐름", {"bold": True, "color": TEXT_LIGHT, "size": 11, "space_after": 4}),
        ("React → Keycloak OIDC 로그인 → Bearer token → /api/* 요청 → Backend가 JWT 검증 + DB 권한 매핑 → 응답", {"color": RGBColor(0xCC, 0xD4, 0xDB), "size": 11}),
    ], line_spacing=1.35)

# ---------- Slide 13: PostgreSQL ERD ----------
def slide_erd():
    s = blank_slide(prs)
    add_header(s, "PostgreSQL — 핵심 ERD", "Gateway 중심. Profile-Channel 분리. Telemetry partition", 13, TOTAL)
    # Build a manual ERD with boxes + connectors
    # Centered layout
    boxes = {
        "companies":         (Inches(0.7),  Inches(1.3),  AMBER),
        "sites":             (Inches(0.7),  Inches(2.4),  AMBER),
        "users":             (Inches(0.7),  Inches(3.6),  TEAL),
        "user_company_roles":(Inches(0.7),  Inches(4.6),  TEAL),
        "user_gateway_perm": (Inches(0.7),  Inches(5.6),  TEAL),
        "gateway_profiles":  (Inches(4.5),  Inches(1.2),  STEEL),
        "gateways":          (Inches(4.5),  Inches(2.4),  NAVY),
        "sensor_profiles":   (Inches(4.5),  Inches(3.7),  STEEL),
        "sensor_channels":   (Inches(4.5),  Inches(4.85), NAVY),
        "actuator_profiles": (Inches(4.5),  Inches(6.0),  STEEL),
        "actuator_channels": (Inches(8.5),  Inches(6.0),  NAVY),
        "gateway_configs":   (Inches(8.5),  Inches(1.2),  GOLD),
        "telemetry":         (Inches(8.5),  Inches(2.4),  GREEN),
        "telemetry_latest":  (Inches(8.5),  Inches(3.5),  GREEN),
        "alarm_rules":       (Inches(8.5),  Inches(4.6),  RED),
        "commands":          (Inches(11),   Inches(2.4),  AMBER),
        "audit_logs":        (Inches(11),   Inches(3.5),  STEEL),
        "bulk_jobs":         (Inches(11),   Inches(4.6),  STEEL),
    }
    box_w = Inches(2.3); box_h = Inches(0.55)
    # Draw boxes
    coords = {}
    for name, (x, y, color) in boxes.items():
        b = add_round_rect(s, x, y, box_w, box_h, color)
        set_text(b, name, size=10.5, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font=FONT_MONO)
        coords[name] = (x, y, x + box_w, y + box_h)

    # Helper: connector between boxes (right-of A → left-of B; or other sides)
    def conn(a, b, side_a="right", side_b="left", color=MUTED, w=0.75, dashed=False):
        ax1, ay1, ax2, ay2 = coords[a]
        bx1, by1, bx2, by2 = coords[b]
        if side_a == "right":   sx = ax2; sy = (ay1 + ay2)/2
        elif side_a == "left":  sx = ax1; sy = (ay1 + ay2)/2
        elif side_a == "bottom":sx = (ax1+ax2)/2; sy = ay2
        else: sx = (ax1+ax2)/2; sy = ay1
        if side_b == "left":    ex = bx1; ey = (by1+by2)/2
        elif side_b == "right": ex = bx2; ey = (by1+by2)/2
        elif side_b == "top":   ex = (bx1+bx2)/2; ey = by1
        else: ex = (bx1+bx2)/2; ey = by2
        add_connector(s, sx, sy, ex, ey, color=color, width=w, arrow=False, dashed=dashed)

    conn("companies", "sites", "bottom", "top")
    conn("sites", "gateways")
    conn("gateway_profiles", "gateways", "bottom", "top")
    conn("gateways", "sensor_channels", "bottom", "top")
    conn("sensor_profiles", "sensor_channels", "bottom", "top")
    conn("gateways", "actuator_channels", "right", "left")
    conn("actuator_profiles", "actuator_channels")
    conn("gateways", "telemetry")
    conn("gateways", "telemetry_latest")
    conn("gateways", "gateway_configs", "right", "left")
    conn("gateways", "commands", "right", "left")
    conn("gateways", "alarm_rules", "right", "left")
    conn("users", "user_company_roles", "bottom", "top")
    conn("users", "user_gateway_perm", "bottom", "top")
    conn("user_gateway_perm", "gateways", "right", "left", dashed=True)

    # Legend
    add_textbox(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.25),
                "■ Identity/Auth   ■ Asset(Gateway)   ■ Profile   ■ Telemetry   ■ Config/Command   ■ Alarm/Audit",
                size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ---------- Slide 14: 핵심 테이블 SQL ----------
def slide_core_tables():
    s = blank_slide(prs)
    add_header(s, "핵심 테이블 (Phase 2 우선)", "companies · sites · users · gateways  ─ DDL 발췌", 14, TOTAL)
    sql_left = """CREATE TABLE companies (
    id UUID PRIMARY KEY,
    name TEXT NOT NULL,
    company_type TEXT NOT NULL,
        -- 'management' | 'customer'
    status TEXT NOT NULL DEFAULT 'active',
    created_at TIMESTAMPTZ DEFAULT now(),
    updated_at TIMESTAMPTZ DEFAULT now()
);

CREATE TABLE sites (
    id UUID PRIMARY KEY,
    company_id UUID NOT NULL
        REFERENCES companies(id),
    name TEXT NOT NULL,
    address TEXT,
    latitude  DOUBLE PRECISION,
    longitude DOUBLE PRECISION,
    created_at TIMESTAMPTZ DEFAULT now()
);

CREATE TABLE users (
    id UUID PRIMARY KEY,
    keycloak_user_id TEXT UNIQUE NOT NULL,
    email TEXT UNIQUE NOT NULL,
    name TEXT NOT NULL,
    status TEXT NOT NULL DEFAULT 'active',
    created_at TIMESTAMPTZ DEFAULT now()
);"""
    sql_right = """CREATE TABLE gateways (
    id UUID PRIMARY KEY,
    serial_number TEXT UNIQUE NOT NULL,
    name TEXT NOT NULL,
    company_id UUID NOT NULL
        REFERENCES companies(id),
    site_id UUID REFERENCES sites(id),
    gateway_profile_id UUID
        REFERENCES gateway_profiles(id),
    status TEXT NOT NULL DEFAULT 'offline',
    firmware_version TEXT,
    app_version TEXT,
    config_version INTEGER DEFAULT 0,
    last_seen_at TIMESTAMPTZ,
    registered_at TIMESTAMPTZ DEFAULT now()
);

CREATE TABLE user_gateway_permissions (
    id UUID PRIMARY KEY,
    user_id UUID NOT NULL REFERENCES users(id),
    gateway_id UUID NOT NULL,
    permission TEXT NOT NULL,
        -- view | control | configure | maintain | admin
    UNIQUE(user_id, gateway_id, permission)
);"""
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(5.95), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.65), Inches(1.18), Inches(5.8), Inches(0.3),
                "companies · sites · users", size=11, bold=True, color=AMBER, font=FONT_MONO)
    add_textbox(s, Inches(0.65), Inches(1.55), Inches(5.8), Inches(5.4),
                sql_left, size=10, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)
    add_round_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(5.95), DARK, line_color=AMBER)
    add_textbox(s, Inches(6.93), Inches(1.18), Inches(5.7), Inches(0.3),
                "gateways · user_gateway_permissions", size=11, bold=True, color=AMBER, font=FONT_MONO)
    add_textbox(s, Inches(6.93), Inches(1.55), Inches(5.8), Inches(5.4),
                sql_right, size=10, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

# ---------- Slide 15: Sensor Profile / Channel ----------
def slide_sensor_profile():
    s = blank_slide(prs)
    add_header(s, "Sensor Profile / Sensor Channel 분리", "센서 종류는 코드가 아니라 데이터로 — 새 센서 추가 시 코드 수정 0", 15, TOTAL)
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(0.5), STEEL)
    add_textbox(s, Inches(0.7), Inches(1.13), Inches(5.5), Inches(0.42),
                "sensor_profiles  (모델 정의)", size=13, bold=True, color=TEXT_LIGHT)
    add_multiline(s, Inches(0.7), Inches(1.7), Inches(5.7), Inches(2.0), [
        ("벤더 / 모델 / 통신 프로토콜", {"size": 11}),
        ("측정 항목 (key, display_name, unit, data_type)", {"size": 11}),
        ("scale / offset / min / max / quality", {"size": 11}),
        ("register map (Modbus function_code, register, length)", {"size": 11}),
        ("default_polling_interval_sec", {"size": 11}),
        ("visualization (line_chart / gauge / map)", {"size": 11}),
    ], bullet=True, line_spacing=1.4)
    add_textbox(s, Inches(0.7), Inches(3.85), Inches(5.7), Inches(0.3),
                "예: TH-RS485-01 (Modbus RTU 온습도)", size=11, bold=True, color=AMBER)
    code1 = """{
  "name": "RS485 Temperature Humidity Sensor",
  "vendor": "Generic", "model": "TH-RS485-01",
  "protocol": "modbus_rtu", "interface_type": "rs485",
  "default_polling_interval_sec": 10,
  "measurements": [
    {"key": "temperature", "unit": "degC",
     "modbus": {"function_code": 3, "register": 0, "length": 1},
     "scale": 0.1, "visualization": "line_chart"},
    {"key": "humidity", "unit": "%",
     "modbus": {"function_code": 3, "register": 1, "length": 1},
     "scale": 0.1, "visualization": "line_chart"}
  ]
}"""
    add_round_rect(s, Inches(0.7), Inches(4.2), Inches(5.7), Inches(2.7), DARK)
    add_textbox(s, Inches(0.85), Inches(4.3), Inches(5.5), Inches(2.55), code1,
                size=8.5, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

    # Right: sensor_channels
    add_round_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(0.5), NAVY)
    add_textbox(s, Inches(6.95), Inches(1.13), Inches(5.7), Inches(0.42),
                "sensor_channels  (Gateway에 실제 연결된 인스턴스)",
                size=13, bold=True, color=TEXT_LIGHT)
    add_multiline(s, Inches(6.95), Inches(1.7), Inches(5.7), Inches(2.0), [
        ("어떤 Gateway에 (gateway_id)", {"size": 11}),
        ("어떤 Sensor Profile을 (sensor_profile_id)", {"size": 11}),
        ("어떤 인터페이스/포트로 (interface_name)", {"size": 11}),
        ("Modbus slave_id / address", {"size": 11}),
        ("polling_interval_sec / enabled / display_name", {"size": 11}),
        ("override 가능한 config JSONB", {"size": 11}),
    ], bullet=True, line_spacing=1.4)
    add_textbox(s, Inches(6.95), Inches(3.85), Inches(5.7), Inches(0.3),
                "예: GW-000001의 1번 RS-485에 슬레이브 1로 연결",
                size=11, bold=True, color=AMBER)
    code2 = """{
  "gateway_id":         "GW-000001",
  "sensor_channel_id":  "sensor-01",
  "sensor_profile_id":  "profile-rs485-temp-humi-001",
  "display_name":       "1번 온습도 센서",
  "interface":          "rs485_1",
  "protocol":           "modbus_rtu",
  "slave_id":           1,
  "polling_interval_sec": 10,
  "enabled":            true
}"""
    add_round_rect(s, Inches(6.95), Inches(4.2), Inches(5.7), Inches(2.7), DARK)
    add_textbox(s, Inches(7.1), Inches(4.3), Inches(5.5), Inches(2.55), code2,
                size=8.5, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

# ---------- Slide 16: Telemetry 저장 모델 ----------
def slide_telemetry():
    s = blank_slide(prs)
    add_header(s, "Telemetry 저장 모델", "월별 partition + telemetry_latest upsert (대시보드용)", 16, TOTAL)
    sql = """CREATE TABLE telemetry (
    id BIGSERIAL,
    company_id UUID NOT NULL,
    site_id UUID,
    gateway_id UUID NOT NULL REFERENCES gateways(id),
    sensor_channel_id UUID NOT NULL REFERENCES sensor_channels(id),
    measurement_key TEXT NOT NULL,
    ts TIMESTAMPTZ NOT NULL,
    value_double DOUBLE PRECISION,
    value_text TEXT, value_bool BOOLEAN, value_json JSONB,
    unit TEXT, quality TEXT, raw JSONB,
    PRIMARY KEY (id, ts)
) PARTITION BY RANGE (ts);

CREATE TABLE telemetry_2026_05 PARTITION OF telemetry
    FOR VALUES FROM ('2026-05-01') TO ('2026-06-01');

CREATE TABLE telemetry_latest (
    gateway_id UUID, sensor_channel_id UUID,
    measurement_key TEXT, ts TIMESTAMPTZ,
    value_double DOUBLE PRECISION,
    value_text TEXT, value_bool BOOLEAN, value_json JSONB,
    unit TEXT, quality TEXT,
    PRIMARY KEY (gateway_id, sensor_channel_id, measurement_key)
);"""
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.95), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.65), Inches(1.15), Inches(7.2), Inches(0.3),
                "DDL — telemetry · telemetry_latest", size=11, bold=True, color=AMBER, font=FONT_MONO)
    add_textbox(s, Inches(0.65), Inches(1.5), Inches(7.2), Inches(5.5), sql,
                size=10, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

    # Right: explanation
    add_round_rect(s, Inches(8.2), Inches(1.1), Inches(4.65), Inches(5.95), SURFACE, line_color=LINE)
    add_textbox(s, Inches(8.4), Inches(1.2), Inches(4.3), Inches(0.4),
                "설계 의도", size=13, bold=True, color=NAVY)
    add_multiline(s, Inches(8.4), Inches(1.6), Inches(4.3), Inches(5.4), [
        ("센서마다 컬럼이 다르므로 measurement_key + value_* 다중 컬럼 패턴", {"size": 10.5}),
        ("월별 partition으로 retention 정책 (drop partition)", {"size": 10.5}),
        ("대시보드는 telemetry_latest 단일 테이블 → O(1) 조회", {"size": 10.5}),
        ("Worker에서 INSERT + UPSERT 동시 수행", {"size": 10.5}),
        ("raw JSONB는 디버깅용 원본 메시지 보존 (옵션)", {"size": 10.5}),
        ("향후 TimescaleDB Community 기능 검토 가능 (License 주의)", {"size": 10, "color": MUTED, "italic": True}),
    ], bullet=True, line_spacing=1.4)

    add_textbox(s, Inches(8.4), Inches(5.7), Inches(4.3), Inches(0.3),
                "예시 데이터", size=11, bold=True, color=AMBER)
    sample = """| GW-001 | sensor-01 | temperature | 24.7 | degC |
| GW-001 | sensor-01 | humidity    | 61.2 | %    |
| GW-002 | sensor-02 | tilt_x      |  3.2 | deg  |"""
    add_textbox(s, Inches(8.4), Inches(6.0), Inches(4.3), Inches(1.0), sample,
                size=9, color=TEXT_DARK, font=FONT_MONO)

# ---------- Slide 17: Gateway Config 버전 ----------
def slide_gw_config():
    s = blank_slide(prs)
    add_header(s, "Gateway Config 버전 관리", "desired_config (서버) ↔ reported_config (Gateway) — 양방향 추적", 17, TOTAL)
    # Flow diagram on top
    steps = [
        ("Web Portal", "관리자 변경", AMBER),
        ("Backend DB", "config 저장", NAVY),
        ("Config Gen", "version + hash", STEEL),
        ("MQTT publish", "gw/{id}/config/desired", TEAL),
        ("Gateway", "설정 적용", AMBER),
        ("MQTT report", "gw/{id}/config/reported", TEAL),
    ]
    sx = Inches(0.5); sy = Inches(1.15); bw = Inches(2.0); bh = Inches(0.85); gap = Inches(0.07)
    for i, (n, d, c) in enumerate(steps):
        x = sx + i * (bw + gap)
        b = add_round_rect(s, x, sy, bw, bh, c)
        add_textbox(s, x, sy + Inches(0.13), bw, Inches(0.32), n,
                    size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(s, x, sy + Inches(0.45), bw, Inches(0.35), d,
                    size=9, color=RGBColor(0xEC, 0xF0, 0xF1), align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_connector(s, x + bw, sy + Inches(0.4), x + bw + gap, sy + Inches(0.4),
                          color=MUTED, width=1.25)

    # Below: SQL + JSON example
    sql = """CREATE TABLE gateway_configs (
    id UUID PRIMARY KEY,
    gateway_id UUID NOT NULL REFERENCES gateways(id),
    config_version INTEGER NOT NULL,
    config_hash TEXT NOT NULL,
    desired_config JSONB NOT NULL,
    status TEXT NOT NULL DEFAULT 'pending',
    created_by UUID REFERENCES users(id),
    created_at TIMESTAMPTZ DEFAULT now(),
    applied_at TIMESTAMPTZ,
    UNIQUE(gateway_id, config_version)
);

CREATE TABLE gateway_config_history (
    id UUID PRIMARY KEY,
    gateway_id UUID NOT NULL REFERENCES gateways(id),
    config_version INTEGER NOT NULL,
    config_snapshot JSONB NOT NULL,
    change_reason TEXT,
    changed_by UUID REFERENCES users(id),
    changed_at TIMESTAMPTZ DEFAULT now()
);"""
    add_round_rect(s, Inches(0.5), Inches(2.25), Inches(6.3), Inches(4.8), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.65), Inches(2.30), Inches(6), Inches(0.3),
                "DDL — gateway_configs · history", size=10.5, bold=True, color=AMBER, font=FONT_MONO)
    add_textbox(s, Inches(0.65), Inches(2.65), Inches(6), Inches(4.3), sql,
                size=9.5, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)
    desired = """// desired (서버)
{
  "gateway_id": "GW-000001",
  "config_version": 12,
  "config_hash": "a83f2e9d",
  "interfaces": [...],
  "sensors":    [...],
  "actuators":  [...],
  "rules":      [...]
}

// reported (Gateway)
{
  "gateway_id": "GW-000001",
  "applied_config_version": 12,
  "config_hash": "a83f2e9d",
  "status": "applied",
  "applied_at": "2026-05-02T12:00:00Z",
  "errors": []
}"""
    add_round_rect(s, Inches(7.0), Inches(2.25), Inches(5.85), Inches(4.8), DARK, line_color=AMBER)
    add_textbox(s, Inches(7.15), Inches(2.30), Inches(5.6), Inches(0.3),
                "JSON — desired vs reported", size=10.5, bold=True, color=AMBER, font=FONT_MONO)
    add_textbox(s, Inches(7.15), Inches(2.65), Inches(5.6), Inches(4.3), desired,
                size=9.5, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

# ---------- Slide 18: MQTT Topic 설계 ----------
def slide_mqtt_topic():
    s = blank_slide(prs)
    add_header(s, "MQTT Topic 설계", "사용자 기준 ❌  ·  Gateway 기준 ✅  (사용자-Gateway 관계는 Backend DB에서)", 18, TOTAL)

    # Bad vs Good
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(1.4), SURFACE, line_color=RED)
    add_textbox(s, Inches(0.7), Inches(1.18), Inches(5.7), Inches(0.3),
                "❌ 비권장", size=11, bold=True, color=RED)
    add_textbox(s, Inches(0.7), Inches(1.5), Inches(5.7), Inches(0.4),
                "user/{userId}/gateway/{gatewayId}/telemetry",
                size=14, color=TEXT_DARK, font=FONT_MONO)
    add_textbox(s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(0.5),
                "사용자 변경/이전 시 topic 재구성 필요. ACL 복잡.",
                size=10, color=MUTED)

    add_round_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(1.4), SURFACE, line_color=GREEN)
    add_textbox(s, Inches(6.95), Inches(1.18), Inches(5.7), Inches(0.3),
                "✅ 권장", size=11, bold=True, color=GREEN)
    add_textbox(s, Inches(6.95), Inches(1.5), Inches(5.7), Inches(0.4),
                "gw/{gatewayId}/telemetry",
                size=14, color=TEXT_DARK, font=FONT_MONO)
    add_textbox(s, Inches(6.95), Inches(1.95), Inches(5.7), Inches(0.5),
                "Gateway 단일 식별자 기준. ACL 단순. 사용자 매핑은 DB.",
                size=10, color=MUTED)

    # Topic table — publish vs subscribe
    add_round_rect(s, Inches(0.5), Inches(2.65), Inches(6.1), Inches(4.4), SURFACE, line_color=LINE)
    add_rect(s, Inches(0.5), Inches(2.65), Inches(6.1), Inches(0.45), NAVY)
    add_textbox(s, Inches(0.7), Inches(2.7), Inches(5.5), Inches(0.4),
                "Gateway → Server (publish)", size=12, bold=True, color=TEXT_LIGHT)
    pub_topics = [
        ("gw/{id}/telemetry",       "센서 측정값"),
        ("gw/{id}/state",           "전체 상태 (cpu·mem·net)"),
        ("gw/{id}/heartbeat",       "주기 alive 신호"),
        ("gw/{id}/event",           "에러·경보·이벤트"),
        ("gw/{id}/config/reported", "현재 적용된 config"),
        ("gw/{id}/command/response","명령 실행 결과"),
        ("gw/{id}/ota/status",      "OTA 진행 상황"),
    ]
    yy = Inches(3.18)
    for t, d in pub_topics:
        add_textbox(s, Inches(0.7), yy, Inches(3.1), Inches(0.3), t,
                    size=10, color=AMBER, font=FONT_MONO, bold=True)
        add_textbox(s, Inches(3.85), yy, Inches(2.7), Inches(0.3), d,
                    size=10, color=TEXT_DARK)
        yy += Inches(0.5)

    add_round_rect(s, Inches(6.78), Inches(2.65), Inches(6.05), Inches(4.4), SURFACE, line_color=LINE)
    add_rect(s, Inches(6.78), Inches(2.65), Inches(6.05), Inches(0.45), STEEL)
    add_textbox(s, Inches(6.95), Inches(2.7), Inches(5.5), Inches(0.4),
                "Server → Gateway (subscribe)", size=12, bold=True, color=TEXT_LIGHT)
    sub_topics = [
        ("gw/{id}/config/desired",  "서버가 발행하는 설정"),
        ("gw/{id}/command/request", "원격 제어 명령"),
        ("gw/{id}/ota/request",     "OTA 작업 지시"),
    ]
    yy = Inches(3.2)
    for t, d in sub_topics:
        add_textbox(s, Inches(6.95), yy, Inches(3.1), Inches(0.3), t,
                    size=10, color=AMBER, font=FONT_MONO, bold=True)
        add_textbox(s, Inches(10.0), yy, Inches(2.7), Inches(0.3), d,
                    size=10, color=TEXT_DARK)
        yy += Inches(0.5)
    # ACL note
    add_round_rect(s, Inches(6.95), Inches(5.0), Inches(5.7), Inches(2.0), DARK, line_color=AMBER)
    add_textbox(s, Inches(7.1), Inches(5.05), Inches(5.4), Inches(0.3),
                "VerneMQ ACL 정책 (Phase 7)", size=11, bold=True, color=AMBER)
    add_multiline(s, Inches(7.1), Inches(5.4), Inches(5.4), Inches(1.5), [
        ("Gateway 계정은 자기 gw/{own_id}/* 만 publish/subscribe 가능", {"size": 10, "color": TEXT_LIGHT}),
        ("Backend 계정은 모든 topic publish/subscribe 가능", {"size": 10, "color": TEXT_LIGHT}),
        ("Phase 1: password file (gateway/admin 1개씩)", {"size": 10, "color": TEXT_LIGHT}),
        ("Phase 7: X.509 client cert + serial → topic ACL 매핑", {"size": 10, "color": TEXT_LIGHT}),
    ], bullet=True, line_spacing=1.35)

# ---------- Slide 19: MQTT Payload ----------
def slide_mqtt_payload():
    s = blank_slide(prs)
    add_header(s, "MQTT Payload 예시", "telemetry · gateway state — 일관된 envelope (timestamp · message_id)", 19, TOTAL)

    tel = """{
  "message_id": "msg-20260502-000001",
  "gateway_id": "GW-000001",
  "timestamp":  "2026-05-02T12:00:00Z",
  "values": [
    {
      "sensor_channel_id": "sensor-01",
      "measurement_key": "temperature",
      "value": 24.7, "unit": "degC", "quality": "good"
    },
    {
      "sensor_channel_id": "sensor-01",
      "measurement_key": "humidity",
      "value": 61.2, "unit": "%", "quality": "good"
    }
  ]
}"""
    state = """{
  "gateway_id": "GW-000001",
  "timestamp":  "2026-05-02T12:00:00Z",
  "status":     "online",
  "app_version":      "1.0.3",
  "firmware_version": "2026.05.02",
  "config_version":   12,
  "uptime_sec":       82344,
  "cpu_temp":         51.3,
  "memory_usage_percent": 38.7,
  "disk_usage_percent":   42.1,
  "network": {
    "ethernet": true,
    "ip":       "192.168.0.25",
    "rtt_ms":   34
  }
}"""
    add_textbox(s, Inches(0.5), Inches(1.05), Inches(6.1), Inches(0.32),
                "Telemetry Payload  →  gw/{id}/telemetry",
                size=12, bold=True, color=AMBER, font=FONT_MONO)
    add_round_rect(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(5.5), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.65), Inches(1.5), Inches(5.85), Inches(5.3), tel,
                size=11, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

    add_textbox(s, Inches(6.78), Inches(1.05), Inches(6.05), Inches(0.32),
                "Gateway State  →  gw/{id}/state",
                size=12, bold=True, color=AMBER, font=FONT_MONO)
    add_round_rect(s, Inches(6.78), Inches(1.4), Inches(6.05), Inches(5.5), DARK, line_color=AMBER)
    add_textbox(s, Inches(6.93), Inches(1.5), Inches(5.85), Inches(5.3), state,
                size=11, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

    add_textbox(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.25),
                "공통 envelope: message_id (idempotency) + gateway_id + timestamp(UTC ISO8601). 모든 메시지 QoS 1 권장.",
                size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ---------- Slide 20: 원격 제어 시퀀스 ----------
def slide_command_flow():
    s = blank_slide(prs)
    add_header(s, "원격 제어 명령 흐름", "User → Backend → MQTT → Gateway → Local Safety → Actuator (양방향 ACK)", 20, TOTAL)
    # Sequence diagram with vertical lifelines
    actors = [("User\nWeb Portal", AMBER), ("Backend\nAPI", NAVY), ("PostgreSQL", GREEN),
              ("VerneMQ", TEAL), ("CM4\nGateway", STEEL), ("Relay\n/Valve", RED)]
    n = len(actors)
    band_x = Inches(0.5); band_w = Inches(12.3); band_top = Inches(1.05)
    col_w = band_w / n
    # Actor headers
    for i, (label, color) in enumerate(actors):
        x = band_x + i * col_w + Inches(0.15)
        bw = col_w - Inches(0.3)
        add_round_rect(s, x, band_top, bw, Inches(0.7), color)
        add_textbox(s, x, band_top + Inches(0.05), bw, Inches(0.6), label,
                    size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # lifeline
        cx = x + bw / 2
        add_connector(s, cx, band_top + Inches(0.7), cx, Inches(7.0),
                      color=LINE, width=0.75, arrow=False, dashed=True)

    def col_x(idx):
        x = band_x + idx * col_w + Inches(0.15)
        bw = col_w - Inches(0.3)
        return x + bw / 2

    # Messages
    msgs = [
        (0, 1, "1) 릴레이 ON 요청 (POST /api/.../commands)"),
        (1, 1, "2) 권한 검사 (User × Gateway)"),
        (1, 2, "3) command 생성 (status=pending)"),
        (1, 3, "4) MQTT publish  gw/{id}/command/request"),
        (3, 4, "5) 명령 전달"),
        (4, 4, "6) Local Safety Rule (interlock · max_on)"),
        (4, 5, "7) GPIO 제어"),
        (5, 4, "8) 실행 결과"),
        (4, 3, "9) MQTT publish  gw/{id}/command/response"),
        (3, 1, "10) Worker가 수신 → status 업데이트"),
        (1, 0, "11) 응답 표시 (executed / failed / timeout)"),
    ]
    y = Inches(2.0)
    for src, dst, label in msgs:
        x1 = col_x(src); x2 = col_x(dst)
        if src == dst:
            # self-loop as small box
            box = add_round_rect(s, x1 - Inches(0.6), y - Inches(0.05), Inches(1.2), Inches(0.3), AMBER)
            set_text(box, "self", size=8, color=DARK, align=PP_ALIGN.CENTER)
        else:
            add_connector(s, x1, y, x2, y, color=NAVY, width=1.25)
        # label
        lx = min(x1, x2) + Inches(0.05)
        lw = abs(x2 - x1) + Inches(2.0)
        add_textbox(s, lx, y - Inches(0.3), lw, Inches(0.28), label,
                    size=9.5, color=TEXT_DARK, anchor=MSO_ANCHOR.BOTTOM)
        y += Inches(0.45)

# ---------- Slide 21: 명령 안전 조건 + Alarm/Rule ----------
def slide_command_safety():
    s = blank_slide(prs)
    add_header(s, "명령 안전 조건 + Alarm/자동 제어 Rule", "원격 제어가 사고로 이어지지 않도록 — 8가지 보호장치", 21, TOTAL)
    # Left: command safety
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(0.5), RED)
    add_textbox(s, Inches(0.7), Inches(1.13), Inches(5.7), Inches(0.42),
                "명령 안전 조건 (8)", size=13, bold=True, color=TEXT_LIGHT)
    safeties = [
        ("command_id", "중복 실행 방지 (idempotency key)"),
        ("expires_at", "오래된 명령 자동 폐기"),
        ("timeout_ms", "지연 명령 실패 처리"),
        ("require_ack", "실행 결과 필수 확인"),
        ("local_safety_check", "현장 조건 위반 시 거부"),
        ("audit_log", "사용자·시간·대상·결과 보관"),
        ("fail_safe", "장애 시 안전 상태로 전환"),
        ("manual_override", "현장 수동 제어 우선"),
    ]
    yy = Inches(1.7)
    for k, v in safeties:
        add_round_rect(s, Inches(0.7), yy, Inches(5.7), Inches(0.6), SUBTLE, line_color=LINE)
        add_textbox(s, Inches(0.85), yy + Inches(0.07), Inches(2.2), Inches(0.45),
                    k, size=10.5, bold=True, color=NAVY, font=FONT_MONO)
        add_textbox(s, Inches(3.1), yy + Inches(0.07), Inches(3.3), Inches(0.45),
                    v, size=10, color=TEXT_DARK)
        yy += Inches(0.65)

    # Right: alarm rule + auto control
    add_round_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(2.85), SURFACE, line_color=LINE)
    add_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(0.5), AMBER)
    add_textbox(s, Inches(6.95), Inches(1.13), Inches(5.7), Inches(0.42),
                "Alarm Rule", size=13, bold=True, color=DARK)
    alarm = """{
  "gateway_id":  "GW-000001",
  "sensor_channel_id": "sensor-01",
  "measurement_key": "temperature",
  "condition": ">", "threshold": 35.0,
  "duration_sec": 60,
  "severity": "warning",
  "action": "notify"
}"""
    add_textbox(s, Inches(6.95), Inches(1.7), Inches(5.7), Inches(2.2), alarm,
                size=10, color=TEXT_DARK, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

    add_round_rect(s, Inches(6.78), Inches(4.1), Inches(6.05), Inches(2.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(6.78), Inches(4.1), Inches(6.05), Inches(0.5), TEAL)
    add_textbox(s, Inches(6.95), Inches(4.13), Inches(5.7), Inches(0.42),
                "자동 제어 Rule (서버 + Gateway 양쪽 배포)",
                size=13, bold=True, color=TEXT_LIGHT)
    ctrl = """{
  "rule_id": "control-rule-001",
  "if": {
    "sensor_channel_id": "sensor-02",
    "measurement_key": "water_level",
    "condition": "<", "threshold": 20
  },
  "then": {
    "actuator_channel_id": "relay-01",
    "action": "ON"
  },
  "safety": { "max_on_duration_sec": 300 }
}"""
    add_textbox(s, Inches(6.95), Inches(4.7), Inches(5.7), Inches(2.3), ctrl,
                size=10, color=TEXT_DARK, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

# ---------- Slide 22: Web Portal 화면 ----------
def slide_web_portal():
    s = blank_slide(prs)
    add_header(s, "Web Portal 화면 구성", "일반 사용자 · Gateway 상세 · 관리자 — 권한별 다른 메뉴", 22, TOTAL)
    cols = [
        ("일반 사용자", AMBER, [
            "내 Gateway 목록",
            "  ─ 이름·위치·온오프라인",
            "  ─ 알람 상태",
            "  ─ 주요 센서 최신값",
            "  ─ 제어 가능 액추에이터",
            "  ─ 최근 이벤트",
        ]),
        ("Gateway 상세", NAVY, [
            "기본 정보 / 네트워크 상태",
            "센서 채널 목록",
            "최신 센서값 (telemetry_latest)",
            "시계열 그래프 (ECharts)",
            "제어 채널 (Toggle/Slider)",
            "알람 이력 / 명령 이력",
            "설정 버전 / 유지보수 로그",
        ]),
        ("관리자", STEEL, [
            "사용자·고객사·현장 관리",
            "Gateway 등록 / 소유권 할당",
            "Gateway Profile 관리",
            "Sensor Profile 관리",
            "Sensor/Actuator Channel 설정",
            "Gateway Template 관리",
            "Bulk Operation / OTA",
            "Audit Log 조회",
        ]),
    ]
    cw = Inches(4.15); ch = Inches(5.95)
    for i, (name, color, items) in enumerate(cols):
        x = Inches(0.5 + i * 4.27)
        y = Inches(1.1)
        add_round_rect(s, x, y, cw, ch, SURFACE, line_color=LINE)
        add_rect(s, x, y, cw, Inches(0.6), color)
        add_textbox(s, x + Inches(0.2), y + Inches(0.1), cw, Inches(0.4),
                    name, size=14, bold=True, color=TEXT_LIGHT)
        bullets = [(t, {"size": 11, "mono": False, "nobullet": t.startswith("  ─")}) for t in items]
        add_multiline(s, x + Inches(0.25), y + Inches(0.75), cw - Inches(0.4), ch - Inches(0.85),
                      bullets, bullet=True, line_spacing=1.4)

# ---------- Slide 23: Sensor Wizard ----------
def slide_sensor_wizard():
    s = blank_slide(prs)
    add_header(s, "Sensor 추가 Wizard", "Gateway별 센서가 다양 → 관리자 UI를 8단계 마법사로", 23, TOTAL)
    steps = [
        ("1", "Gateway 선택", "어떤 장비에 추가하나"),
        ("2", "인터페이스 선택", "RS-485 #1/2 · AI · DI · I2C · UART"),
        ("3", "Sensor Profile 선택", "온습도 · 미세먼지 · 기울기 · 수위 · pH ..."),
        ("4", "통신 설정", "Modbus slave_id · baudrate · parity"),
        ("5", "측정 주기", "polling_interval_sec"),
        ("6", "표시 이름", "사용자에게 보여줄 이름"),
        ("7", "저장", "sensor_channels 테이블에 INSERT"),
        ("8", "Config 배포", "config_version + 1, MQTT desired publish"),
    ]
    y = Inches(1.15)
    cw = Inches(12.3); ch = Inches(0.65)
    for n, t, d in steps:
        add_round_rect(s, Inches(0.5), y, cw, ch, SURFACE, line_color=LINE)
        add_round_rect(s, Inches(0.55), y + Inches(0.07), Inches(0.5), Inches(0.5), AMBER)
        add_textbox(s, Inches(0.55), y + Inches(0.08), Inches(0.5), Inches(0.5),
                    n, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(1.2), y + Inches(0.05), Inches(3.5), Inches(0.55),
                    t, size=12.5, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(4.7), y + Inches(0.05), Inches(8), Inches(0.55),
                    d, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.72)

# ---------- Slide 24: Dashboard 자동 생성 ----------
def slide_dynamic_dashboard():
    s = blank_slide(prs)
    add_header(s, "동적 Dashboard 자동 생성", "Sensor Profile의 visualization hint → Widget 매핑", 24, TOTAL)
    # Mapping table
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(7.5), Inches(0.5), NAVY)
    add_textbox(s, Inches(0.7), Inches(1.13), Inches(7), Inches(0.42),
                "measurement_key  →  Widget 매핑", size=13, bold=True, color=TEXT_LIGHT)
    rows = [
        ("temperature",   "Line Chart + Current Value Card"),
        ("humidity",      "Line Chart + Gauge"),
        ("pressure",      "Line Chart"),
        ("pm2_5",         "Gauge + Line Chart"),
        ("tilt_x, tilt_y","2-axis Tilt View + Line"),
        ("water_level",   "Gauge"),
        ("relay_state",   "Toggle + Status Card"),
        ("valve_state",   "Toggle + Status Card"),
        ("gps",           "Map"),
    ]
    y = Inches(1.7)
    for i, (k, w) in enumerate(rows):
        bg = SURFACE if i % 2 == 0 else SUBTLE
        add_rect(s, Inches(0.5), y, Inches(7.5), Inches(0.45), bg, line_color=LINE, line_w=0.25)
        add_textbox(s, Inches(0.7), y + Inches(0.06), Inches(2.6), Inches(0.32),
                    k, size=11, color=AMBER, font=FONT_MONO, bold=True)
        add_textbox(s, Inches(3.4), y + Inches(0.06), Inches(4.5), Inches(0.32),
                    w, size=10.5, color=TEXT_DARK)
        y += Inches(0.45)

    # Right: example
    add_round_rect(s, Inches(8.2), Inches(1.1), Inches(4.65), Inches(5.95), DARK, line_color=AMBER)
    add_textbox(s, Inches(8.35), Inches(1.18), Inches(4.4), Inches(0.3),
                "Sensor Profile 예", size=11, bold=True, color=AMBER, font=FONT_MONO)
    code = """{
  "key": "temperature",
  "display_name": "Temperature",
  "unit": "degC",
  "visualization": "line_chart",
  "display_group": "environment",
  "order": 1
}"""
    add_textbox(s, Inches(8.35), Inches(1.5), Inches(4.4), Inches(2.2), code,
                size=10, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)
    add_textbox(s, Inches(8.35), Inches(3.85), Inches(4.4), Inches(0.4),
                "결과: Frontend가 자동으로 Line Chart를 environment 그룹에 1번째로 배치",
                size=10, color=RGBColor(0xCC, 0xD4, 0xDB), italic=True)
    add_textbox(s, Inches(8.35), Inches(4.5), Inches(4.4), Inches(2.4),
                "장점\n• 새 센서 추가 시 frontend 수정 0\n• Gateway별 다른 대시보드 자동 생성\n• 사용자별 widget 순서 커스텀 가능",
                size=11, color=TEXT_LIGHT)

# ---------- Slide 25: Bulk Operation ----------
def slide_bulk_ops():
    s = blank_slide(prs)
    add_header(s, "Bulk Operation", "Gateway 수가 늘면 일괄 작업이 필수 — 8가지 작업 타입", 25, TOTAL)
    ops = [
        ("Gateway 현장 일괄 배정", "company/site 매핑 변경"),
        ("사용자 권한 일괄 부여",  "특정 group → 다수 Gateway"),
        ("센서 polling 일괄 변경",  "sensor_channels.polling_interval_sec"),
        ("알람 기준 일괄 변경",    "alarm_rules.threshold"),
        ("OTA 일괄 업데이트",       "타겟 필터 + 점진 배포"),
        ("Gateway 일괄 재시작",     "command publish (reboot)"),
        ("설정 일괄 배포",          "Template 기반 desired_config 발행"),
        ("로그 일괄 수집",          "log/upload topic 트리거"),
    ]
    cw = Inches(6.05); ch = Inches(0.65)
    for i, (n, d) in enumerate(ops):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.27)
        y = Inches(1.15 + row * 0.75)
        add_round_rect(s, x, y, cw, ch, SURFACE, line_color=LINE)
        b = add_round_rect(s, x + Inches(0.1), y + Inches(0.1), Inches(0.45), Inches(0.45), AMBER)
        set_text(b, str(i+1), size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.7), y + Inches(0.05), Inches(2.8), Inches(0.55),
                    n, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, x + Inches(3.5), y + Inches(0.05), Inches(2.5), Inches(0.55),
                    d, size=10.5, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom: schema
    add_round_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.7), Inches(4.55), Inches(12), Inches(0.3),
                "DDL — bulk_jobs", size=11, bold=True, color=AMBER, font=FONT_MONO)
    sql = """CREATE TABLE bulk_jobs (
    id UUID PRIMARY KEY,
    job_type TEXT NOT NULL,
    target_filter JSONB NOT NULL,    -- {"company_id": "...", "site_id": "...", "tags": [...]}
    payload JSONB NOT NULL,          -- 작업별 파라미터
    status TEXT NOT NULL DEFAULT 'pending',
    created_by UUID REFERENCES users(id),
    created_at TIMESTAMPTZ DEFAULT now(),
    completed_at TIMESTAMPTZ
);"""
    add_textbox(s, Inches(0.7), Inches(4.9), Inches(12), Inches(2.0), sql,
                size=10.5, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

# ---------- Slide 26: Gateway Agent 구현 ----------
def slide_gateway_agent():
    s = blank_slide(prs)
    add_header(s, "Gateway Agent 구현 계획", "/opt/iot-gateway/ 모듈 구성 + 부팅 흐름 + Local Buffer", 26, TOTAL)

    # Left: directory
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(4.1), Inches(5.95), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.65), Inches(1.18), Inches(3.8), Inches(0.3),
                "/opt/iot-gateway/", size=11, bold=True, color=AMBER, font=FONT_MONO)
    tree = """gateway-agent/      Main loop
sensor-service/     Modbus poll
actuator-service/   GPIO/Relay
rule-engine/        Local Rule
mqtt-client/        TLS pub/sub
local-db/           SQLite
ota-agent/          이미지 검증
health-agent/       heartbeat
config/             설정 파일
logs/               회전 로그"""
    add_textbox(s, Inches(0.65), Inches(1.55), Inches(3.8), Inches(5.4), tree,
                size=11, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

    # Middle: boot flow
    add_round_rect(s, Inches(4.78), Inches(1.1), Inches(4.1), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(4.78), Inches(1.1), Inches(4.1), Inches(0.5), NAVY)
    add_textbox(s, Inches(4.95), Inches(1.13), Inches(3.8), Inches(0.42),
                "부팅 흐름 (14 단계)", size=13, bold=True, color=TEXT_LIGHT)
    boot = [
        "Gateway ID 확인",
        "인증서·설정 파일 확인",
        "네트워크 확인",
        "MQTT Broker 연결",
        "현재 config_version 보고",
        "서버 desired_config 확인",
        "버전 다르면 다운로드",
        "설정 유효성 검사",
        "Sensor driver 구성",
        "Actuator driver 구성",
        "Rule engine 구성",
        "reported_config 전송",
        "센서 polling 시작",
        "주기 heartbeat 전송",
    ]
    yy = Inches(1.7)
    for i, b in enumerate(boot):
        add_textbox(s, Inches(4.95), yy, Inches(0.3), Inches(0.32),
                    f"{i+1:02d}", size=10, color=AMBER, bold=True, font=FONT_MONO)
        add_textbox(s, Inches(5.3), yy, Inches(3.4), Inches(0.32), b, size=10.5, color=TEXT_DARK)
        yy += Inches(0.36)

    # Right: local buffer
    add_round_rect(s, Inches(9.05), Inches(1.1), Inches(3.8), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(9.05), Inches(1.1), Inches(3.8), Inches(0.5), STEEL)
    add_textbox(s, Inches(9.2), Inches(1.13), Inches(3.5), Inches(0.42),
                "Local Buffer 정책 (SQLite)", size=12, bold=True, color=TEXT_LIGHT)
    items = [
        ("저장 대상", "telemetry · event · cmd response"),
        ("재전송 순서", "timestamp 오름차순"),
        ("중복 방지", "message_id 기반"),
        ("보존 기간", "7~30일 (정책 가능)"),
        ("초과 시", "오래된 telemetry부터 삭제"),
        ("우선순위", "event > cmd_resp > telemetry"),
    ]
    yy = Inches(1.7)
    for k, v in items:
        add_textbox(s, Inches(9.2), yy, Inches(3.5), Inches(0.25), k,
                    size=10, bold=True, color=AMBER)
        add_textbox(s, Inches(9.2), yy + Inches(0.25), Inches(3.5), Inches(0.45), v,
                    size=9.5, color=TEXT_DARK)
        yy += Inches(0.78)

# ---------- Slide 27: Safety MCU ----------
def slide_safety():
    s = blank_slide(prs)
    add_header(s, "Gateway Safety 설계", "CM4 Linux + STM32/NXP Safety MCU 이중 구조 (Phase 4-7 도입)", 27, TOTAL)

    # CM4 box
    add_round_rect(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(2.4), STEEL)
    add_textbox(s, Inches(0.7), Inches(1.25), Inches(5.7), Inches(0.4),
                "CM4 Linux Gateway", size=14, bold=True, color=TEXT_LIGHT)
    add_multiline(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(1.8), [
        ("클라우드 통신 (MQTT TLS)", {"size": 11, "color": TEXT_LIGHT}),
        ("데이터 저장 (SQLite local buffer)", {"size": 11, "color": TEXT_LIGHT}),
        ("고수준 명령 처리 / Rule engine", {"size": 11, "color": TEXT_LIGHT}),
        ("OTA / 진단 / 로그 수집", {"size": 11, "color": TEXT_LIGHT}),
    ], bullet=True)

    # Arrow down
    add_connector(s, Inches(3.5), Inches(3.6), Inches(3.5), Inches(4.0), color=NAVY, width=2.5)

    # Safety MCU
    add_round_rect(s, Inches(0.5), Inches(4.05), Inches(6.0), Inches(3.0), RED)
    add_textbox(s, Inches(0.7), Inches(4.15), Inches(5.7), Inches(0.4),
                "STM32 / NXP Safety MCU", size=14, bold=True, color=TEXT_LIGHT)
    add_multiline(s, Inches(0.7), Inches(4.6), Inches(5.6), Inches(2.4), [
        ("릴레이/밸브 직접 제어 (low-level)", {"size": 11, "color": TEXT_LIGHT}),
        ("Local Interlock (센서 조건 위반 시 차단)", {"size": 11, "color": TEXT_LIGHT}),
        ("Watchdog (Linux 행 시 자동 fail-safe)", {"size": 11, "color": TEXT_LIGHT}),
        ("Fail-safe state 강제 전환", {"size": 11, "color": TEXT_LIGHT}),
        ("물리적 비상 정지 입력 처리", {"size": 11, "color": TEXT_LIGHT}),
    ], bullet=True)

    # Right: 8 essential safety features
    add_round_rect(s, Inches(6.78), Inches(1.15), Inches(6.05), Inches(5.9), SURFACE, line_color=LINE)
    add_textbox(s, Inches(6.95), Inches(1.25), Inches(5.7), Inches(0.4),
                "필수 안전 기능 (8)", size=14, bold=True, color=NAVY)
    feats = [
        ("Fail-safe state",    "장애 시 릴레이/밸브 기본 안전 상태"),
        ("Max ON duration",    "릴레이가 너무 오래 켜지지 않도록"),
        ("Command expiry",     "오래된 명령 실행 금지"),
        ("Manual override",    "현장 수동 제어 우선"),
        ("Interlock",          "센서 조건 위반 시 제어 차단"),
        ("Watchdog",           "프로세스/OS 장애 감지"),
        ("Output feedback",    "실제 릴레이 상태 피드백"),
        ("Emergency stop",     "물리적 비상 정지"),
    ]
    yy = Inches(1.75)
    for k, v in feats:
        add_round_rect(s, Inches(6.95), yy, Inches(5.7), Inches(0.6), SUBTLE, line_color=LINE)
        add_textbox(s, Inches(7.1), yy + Inches(0.07), Inches(2.0), Inches(0.45),
                    k, size=11, bold=True, color=RED)
        add_textbox(s, Inches(9.2), yy + Inches(0.07), Inches(3.5), Inches(0.45),
                    v, size=10, color=TEXT_DARK)
        yy += Inches(0.65)

# ---------- Slide 28: API 설계 ----------
def slide_api_design():
    s = blank_slide(prs)
    add_header(s, "API 설계 요약", "Gateway · Sensor · Actuator · Config · Admin — REST 패턴", 28, TOTAL)
    groups = [
        ("Gateway", AMBER, [
            "POST   /api/gateways",
            "GET    /api/gateways",
            "GET    /api/gateways/{id}",
            "PATCH  /api/gateways/{id}",
            "DELETE /api/gateways/{id}",
            "GET    /api/gateways/{id}/state",
            "GET    /api/gateways/{id}/telemetry",
            "GET    /api/gateways/{id}/latest",
            "GET    /api/gateways/{id}/events",
        ]),
        ("Sensor", NAVY, [
            "POST   /api/sensor-profiles",
            "GET    /api/sensor-profiles",
            "PATCH  /api/sensor-profiles/{id}",
            "POST   /api/gateways/{id}/sensor-channels",
            "GET    /api/gateways/{id}/sensor-channels",
            "PATCH  /api/sensor-channels/{id}",
            "DELETE /api/sensor-channels/{id}",
        ]),
        ("Actuator + Command", STEEL, [
            "POST   /api/actuator-profiles",
            "GET    /api/actuator-profiles",
            "POST   /api/gateways/{id}/actuator-channels",
            "GET    /api/gateways/{id}/actuator-channels",
            "PATCH  /api/actuator-channels/{id}",
            "POST   /api/gateways/{id}/commands",
            "GET    /api/commands/{cmd_id}",
        ]),
        ("Config + Admin", TEAL, [
            "POST   /api/gateways/{id}/configs/generate",
            "GET    /api/gateways/{id}/configs",
            "GET    /api/gateways/{id}/configs/latest",
            "POST   /api/gateways/{id}/configs/{ver}/deploy",
            "POST   /api/gateways/{id}/configs/{ver}/rollback",
            "POST   /api/companies   /api/sites",
            "POST   /api/users/{id}/gateway-permissions",
            "GET    /api/audit-logs   /api/bulk-jobs",
        ]),
    ]
    cw = Inches(6.1); ch = Inches(2.95)
    coords = [(Inches(0.5), Inches(1.1)), (Inches(6.78), Inches(1.1)),
              (Inches(0.5), Inches(4.1)), (Inches(6.78), Inches(4.1))]
    for (name, color, items), (x, y) in zip(groups, coords):
        add_round_rect(s, x, y, cw, ch, SURFACE, line_color=LINE)
        add_rect(s, x, y, cw, Inches(0.45), color)
        add_textbox(s, x + Inches(0.15), y + Inches(0.05), cw, Inches(0.35),
                    name, size=12, bold=True, color=TEXT_LIGHT)
        add_multiline(s, x + Inches(0.2), y + Inches(0.55), cw - Inches(0.4), ch - Inches(0.65),
                      [(t, {"mono": True, "size": 9.5, "color": TEXT_DARK}) for t in items],
                      line_spacing=1.25)

# ---------- Slide 29: 보안 설계 ----------
def slide_security():
    s = blank_slide(prs)
    add_header(s, "보안 설계", "장비 보안 + 서버 보안 — TLS · JWT · ACL · Audit 다중 방어선", 29, TOTAL)
    # Two columns
    dev = [
        ("MQTT", "TLS 필수 (Phase 7)"),
        ("Gateway 인증", "Phase 1 ID/PW · Phase 7 X.509"),
        ("Topic ACL", "자기 topic만"),
        ("Private Key", "TPM/Secure Element 권장"),
        ("SSH", "기본 비활성화"),
        ("OTA", "서명 검증 필수"),
        ("Local Config", "config_hash 검증"),
        ("로그", "제어·설정 변경 보관"),
    ]
    srv = [
        ("인증",   "Keycloak OIDC/OAuth2"),
        ("API",    "JWT 검증 (Backend)"),
        ("권한",   "RBAC + ABAC"),
        ("DB 필터","company/site/gateway_id"),
        ("RLS",    "주요 테이블 적용 검토"),
        ("TLS",    "Web · API · MQTT 모두"),
        ("Audit",  "사용자 명령·설정 변경"),
        ("Backup", "정기 + 복구 테스트"),
    ]
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(6.1), Inches(0.5), STEEL)
    add_textbox(s, Inches(0.7), Inches(1.13), Inches(5.7), Inches(0.42),
                "장비 (Gateway) 보안", size=13, bold=True, color=TEXT_LIGHT)
    yy = Inches(1.75)
    for k, v in dev:
        add_round_rect(s, Inches(0.7), yy, Inches(5.7), Inches(0.55), SUBTLE, line_color=LINE)
        add_textbox(s, Inches(0.85), yy + Inches(0.05), Inches(2.0), Inches(0.45),
                    k, size=11, bold=True, color=NAVY)
        add_textbox(s, Inches(2.9), yy + Inches(0.05), Inches(3.5), Inches(0.45),
                    v, size=10, color=TEXT_DARK)
        yy += Inches(0.62)

    add_round_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(6.78), Inches(1.1), Inches(6.05), Inches(0.5), NAVY)
    add_textbox(s, Inches(6.95), Inches(1.13), Inches(5.7), Inches(0.42),
                "서버 보안", size=13, bold=True, color=TEXT_LIGHT)
    yy = Inches(1.75)
    for k, v in srv:
        add_round_rect(s, Inches(6.95), yy, Inches(5.7), Inches(0.55), SUBTLE, line_color=LINE)
        add_textbox(s, Inches(7.1), yy + Inches(0.05), Inches(2.0), Inches(0.45),
                    k, size=11, bold=True, color=NAVY)
        add_textbox(s, Inches(9.15), yy + Inches(0.05), Inches(3.5), Inches(0.45),
                    v, size=10, color=TEXT_DARK)
        yy += Inches(0.62)

# ---------- Slide 30: RLS 검토 ----------
def slide_rls():
    s = blank_slide(prs)
    add_header(s, "PostgreSQL Row Level Security 검토", "사용자별 Gateway 접근 격리 — Phase 6+ 적용 권장", 30, TOTAL)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4),
                "적용 대상 테이블", size=12, bold=True, color=AMBER)
    targets = ["gateways", "sensor_channels", "actuator_channels",
               "telemetry", "telemetry_latest", "commands", "audit_logs"]
    x = Inches(0.5)
    for t in targets:
        b = add_round_rect(s, x, Inches(1.5), Inches(1.65), Inches(0.45), SURFACE, line_color=NAVY)
        set_text(b, t, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_MONO)
        x += Inches(1.78)

    sql = """ALTER TABLE gateways ENABLE ROW LEVEL SECURITY;

CREATE POLICY gateway_access_policy ON gateways
USING (
    id IN (
        SELECT gateway_id
          FROM user_gateway_permissions
         WHERE user_id = current_setting('app.current_user_id')::uuid
    )
);"""
    add_round_rect(s, Inches(0.5), Inches(2.2), Inches(7.5), Inches(3.2), DARK, line_color=AMBER)
    add_textbox(s, Inches(0.65), Inches(2.25), Inches(7.2), Inches(0.3),
                "예시 정책 (PoC)", size=11, bold=True, color=AMBER, font=FONT_MONO)
    add_textbox(s, Inches(0.65), Inches(2.6), Inches(7.2), Inches(2.7), sql,
                size=11, color=TEXT_LIGHT, font=FONT_MONO, anchor=MSO_ANCHOR.TOP)

    add_round_rect(s, Inches(8.2), Inches(2.2), Inches(4.65), Inches(3.2), SURFACE, line_color=LINE)
    add_textbox(s, Inches(8.4), Inches(2.3), Inches(4.3), Inches(0.4),
                "도입 시 고려사항", size=12, bold=True, color=NAVY)
    add_multiline(s, Inches(8.4), Inches(2.7), Inches(4.3), Inches(2.7), [
        ("Backend connection pool에서 SET app.current_user_id 호출", {"size": 10}),
        ("service role 우회 정책 정의 필요 (worker, scheduler)", {"size": 10}),
        ("관리자 권한 BYPASSRLS 별도 role", {"size": 10}),
        ("batch/migration 작업 시 일시 비활성화 고려", {"size": 10}),
        ("Phase 6+ 적용 권장 — 초기에는 application 단 권한으로", {"size": 10, "color": AMBER, "italic": True}),
    ], bullet=True, line_spacing=1.4)

    add_round_rect(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.45), SUBTLE, line_color=LINE)
    add_textbox(s, Inches(0.7), Inches(5.65), Inches(12), Inches(0.35),
                "결론", size=12, bold=True, color=NAVY)
    add_textbox(s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.0),
                "Phase 1-5 동안은 application 레이어에서 권한 enforce. "
                "DB row-level 격리는 다중 테넌트가 본격화되는 Phase 6+에 도입.",
                size=11, color=TEXT_DARK)

# ---------- Slide 31: 운영 관리 ----------
def slide_operations():
    s = blank_slide(prs)
    add_header(s, "운영 관리 — 백업 · 모니터링 · 장애 대응", "Phase 7 제품화 단계의 핵심 — 자체 호스팅이므로 모두 직접 책임", 31, TOTAL)

    # Backup
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(4.1), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(4.1), Inches(0.5), TEAL)
    add_textbox(s, Inches(0.7), Inches(1.13), Inches(3.7), Inches(0.42),
                "백업", size=13, bold=True, color=TEXT_LIGHT)
    backups = [
        ("PostgreSQL", "매일", "pg_dump / pg_basebackup"),
        ("Gateway Config", "매일", "파일/DB 백업"),
        ("펌웨어", "변경 시", "rsync"),
        ("로그", "정책", "압축/보관"),
        ("Keycloak realm", "변경 시", "kc.sh export"),
    ]
    yy = Inches(1.75)
    for k, freq, how in backups:
        add_textbox(s, Inches(0.7), yy, Inches(3.7), Inches(0.25), k,
                    size=11, bold=True, color=AMBER)
        add_textbox(s, Inches(0.7), yy + Inches(0.25), Inches(1.3), Inches(0.25),
                    freq, size=10, color=NAVY, italic=True)
        add_textbox(s, Inches(2.0), yy + Inches(0.25), Inches(2.5), Inches(0.25),
                    how, size=10, color=TEXT_DARK, font=FONT_MONO)
        yy += Inches(0.7)

    # Monitoring
    add_round_rect(s, Inches(4.78), Inches(1.1), Inches(4.0), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(4.78), Inches(1.1), Inches(4.0), Inches(0.5), NAVY)
    add_textbox(s, Inches(4.95), Inches(1.13), Inches(3.7), Inches(0.42),
                "모니터링 지표", size=13, bold=True, color=TEXT_LIGHT)
    metrics = [
        ("Gateway", "online · heartbeat · CPU · mem · disk · temp"),
        ("VerneMQ", "connected · message rate · dropped"),
        ("Backend", "API latency · error rate"),
        ("DB", "connections · slow query · disk"),
        ("Command", "success · timeout · rejected rate"),
        ("Config", "pending · applied · failed"),
        ("OTA", "success · failed · rollback"),
    ]
    yy = Inches(1.75)
    for k, v in metrics:
        add_textbox(s, Inches(4.95), yy, Inches(1.4), Inches(0.7), k,
                    size=11, bold=True, color=AMBER)
        add_textbox(s, Inches(6.4), yy, Inches(2.3), Inches(0.7), v,
                    size=9.5, color=TEXT_DARK)
        yy += Inches(0.7)

    # Failure response
    add_round_rect(s, Inches(8.95), Inches(1.1), Inches(3.9), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(8.95), Inches(1.1), Inches(3.9), Inches(0.5), RED)
    add_textbox(s, Inches(9.12), Inches(1.13), Inches(3.7), Inches(0.42),
                "장애 대응", size=13, bold=True, color=TEXT_LIGHT)
    failures = [
        ("Gateway offline", "알람 + 마지막 상태 표시"),
        ("센서 미수신", "channel = degraded"),
        ("명령 timeout", "command failed"),
        ("config 적용 실패", "이전 config 유지 + 알림"),
        ("DB 용량 증가", "partition retention"),
        ("서버 장애", "백업 복구 절차"),
    ]
    yy = Inches(1.75)
    for k, v in failures:
        add_textbox(s, Inches(9.12), yy, Inches(3.6), Inches(0.25), k,
                    size=11, bold=True, color=AMBER)
        add_textbox(s, Inches(9.12), yy + Inches(0.25), Inches(3.6), Inches(0.5), v,
                    size=10, color=TEXT_DARK)
        yy += Inches(0.85)

# ---------- Slide 32: 개발 로드맵 7단계 ----------
def slide_roadmap():
    s = blank_slide(prs)
    add_header(s, "개발 로드맵 — 7단계", "Phase 1만 우선 1주 sprint로 확정. 나머지는 결과 보고 결정", 32, TOTAL)
    phases = [
        ("Phase 1", "서버 기본 구축", "1주",   "systemd 5종 · Keycloak realm · 코드 0",     AMBER),
        ("Phase 2", "다중 Gateway 권한 모델", "TBD", "companies/sites/users/gateways · Backend skeleton", NAVY),
        ("Phase 3", "Sensor Profile/Channel", "TBD", "Sensor Wizard · Telemetry · 동적 Dashboard",        STEEL),
        ("Phase 4", "Actuator + 원격 제어", "TBD",  "command 흐름 · 안전 조건 · 권한 분리",            TEAL),
        ("Phase 5", "Gateway Config Versioning", "TBD", "desired/reported · rollback",                  GOLD),
        ("Phase 6", "관리 편의성 (Template/Bulk/Alarm)", "TBD", "Template · Bulk · Alarm Rule",           STEEL),
        ("Phase 7", "제품화", "TBD",          "X.509 · ACL · OTA · Backup · Safety MCU · OSS Notice",   RED),
    ]
    y = Inches(1.1)
    for ph, name, dur, items, color in phases:
        add_round_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.78), SURFACE, line_color=LINE)
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.78), color)
        add_textbox(s, Inches(0.85), y + Inches(0.07), Inches(1.5), Inches(0.32),
                    ph, size=13, bold=True, color=NAVY)
        add_textbox(s, Inches(2.4), y + Inches(0.07), Inches(4.5), Inches(0.32),
                    name, size=12, bold=True, color=TEXT_DARK)
        add_textbox(s, Inches(0.85), y + Inches(0.42), Inches(1.5), Inches(0.32),
                    dur, size=10, color=AMBER, bold=True)
        add_textbox(s, Inches(2.4), y + Inches(0.42), Inches(10.0), Inches(0.32),
                    items, size=10.5, color=TEXT_DARK)
        y += Inches(0.85)

# ---------- Slide 33: Phase 1 상세 spec ----------
def slide_phase1_spec():
    s = blank_slide(prs)
    add_header(s, "Phase 1 — 서버 인프라 1주 sprint (확정 spec)",
               "인터뷰 결과 · 코드 0줄 · Solo + Claude Code AI pair", 33, TOTAL)

    # Top: meta strip
    metas = [("환경", "Ubuntu 24.04 · 사내 물리"),
             ("기간", "1주 single sprint"),
             ("코드", "0줄 (순수 인프라)"),
             ("팀",   "Solo + AI pair"),
             ("Verification", "체크리스트 + smoke.sh")]
    x = Inches(0.5)
    for k, v in metas:
        bw = Inches(2.4)
        add_round_rect(s, x, Inches(1.1), bw, Inches(0.7), DARK, line_color=AMBER)
        add_textbox(s, x + Inches(0.1), Inches(1.15), bw, Inches(0.3),
                    k, size=10, color=AMBER, bold=True)
        add_textbox(s, x + Inches(0.1), Inches(1.42), bw - Inches(0.1), Inches(0.35),
                    v, size=10.5, color=TEXT_LIGHT)
        x += bw + Inches(0.05)

    # Scope (left)
    add_round_rect(s, Inches(0.5), Inches(1.95), Inches(6.1), Inches(5.1), SURFACE, line_color=GREEN, line_w=1)
    add_rect(s, Inches(0.5), Inches(1.95), Inches(6.1), Inches(0.45), GREEN)
    add_textbox(s, Inches(0.7), Inches(1.99), Inches(5.7), Inches(0.4),
                "✅ Scope (Phase 1 포함)", size=12, bold=True, color=TEXT_LIGHT)
    scope = [
        "OS · 보안 hardening (ufw · ssh · iot 사용자)",
        "PostgreSQL 16 (apt) · iot_platform DB · keycloak DB · iot_user role",
        "VerneMQ (.deb) · systemd · password file 인증 · 1883 plain",
        "Keycloak (release) · systemd · Postgres backend",
        "    └ realm: iot-platform · 7 roles · test user 1명",
        "Nginx + certbot Let's Encrypt → HTTPS reverse proxy",
        "    └ /auth → keycloak (placeholder /api · /)",
        "/etc/iot-platform/ · .env 템플릿",
        "운영 절차서 (markdown)",
        "scripts/phase1_smoke.sh (자동 점검)",
    ]
    add_multiline(s, Inches(0.7), Inches(2.5), Inches(5.7), Inches(4.5),
                  [(t, {"size": 10.5, "nobullet": t.startswith("    └")}) for t in scope],
                  bullet=True, line_spacing=1.4)

    # Non-goals (right top)
    add_round_rect(s, Inches(6.78), Inches(1.95), Inches(6.05), Inches(2.45), SURFACE, line_color=RED, line_w=1)
    add_rect(s, Inches(6.78), Inches(1.95), Inches(6.05), Inches(0.45), RED)
    add_textbox(s, Inches(6.95), Inches(1.99), Inches(5.7), Inches(0.4),
                "❌ Non-goals (Phase 2+ 이연)", size=12, bold=True, color=TEXT_LIGHT)
    non = [
        "모든 application 코드 (FastAPI · React · Worker · Scheduler)",
        "비즈니스 DB schema (companies · sites · users · gateways · ...)",
        "Gateway 연결 · telemetry · MQTT TLS · X.509",
        "OTA · Safety MCU · Backup automation · Alarm · Bulk",
    ]
    add_multiline(s, Inches(6.95), Inches(2.5), Inches(5.7), Inches(1.85),
                  [(t, {"size": 10.5}) for t in non], bullet=True, line_spacing=1.4)

    # DoD (right bottom)
    add_round_rect(s, Inches(6.78), Inches(4.5), Inches(6.05), Inches(2.55), SURFACE, line_color=NAVY, line_w=1)
    add_rect(s, Inches(6.78), Inches(4.5), Inches(6.05), Inches(0.45), NAVY)
    add_textbox(s, Inches(6.95), Inches(4.55), Inches(5.7), Inches(0.4),
                "🎯 Definition of Done", size=12, bold=True, color=TEXT_LIGHT)
    dod = [
        "systemctl is-active → 4개 서비스 active",
        "https://<도메인>/auth/realms/iot-platform/.well-known/openid-configuration → 200",
        "Keycloak admin 로그인 성공 (test user)",
        "mosquitto_pub/sub → password 인증 후 1883 publish/subscribe 성공",
        "psql -U iot_user -d iot_platform -c \"SELECT 1\" → 성공",
        "scripts/phase1_smoke.sh → 전체 PASS",
    ]
    add_multiline(s, Inches(6.95), Inches(5.05), Inches(5.7), Inches(2.0),
                  [(t, {"size": 10, "mono": True}) for t in dod], bullet=True, line_spacing=1.4)

# ---------- Slide 34: Phase 1 일정 분해 ----------
def slide_phase1_schedule():
    s = blank_slide(prs)
    add_header(s, "Phase 1 일정 분해 (1주 sprint 제안)", "Day 1-7 task breakdown — 인터뷰 합의 spec 기준", 34, TOTAL)
    days = [
        ("Day 1", "OS 베이스 + 보안", [
            "Ubuntu 24.04 LTS 클린 설치",
            "iot 시스템 사용자 생성 · sudo 정책",
            "ufw 포트 정책 (22 · 80 · 443 · 1883)",
            "ssh hardening (key only · port forward 차단)",
            "/opt/iot-platform · /etc/iot-platform · /var/lib/iot-platform 디렉터리"
        ]),
        ("Day 2", "PostgreSQL + VerneMQ", [
            "PostgreSQL 16 apt 설치 · 권한 설정",
            "iot_platform DB + keycloak DB 생성",
            "iot_user role + 권한 grant",
            "VerneMQ .deb 설치 · systemd",
            "vmq_passwd로 gateway/admin 계정 생성"
        ]),
        ("Day 3", "Keycloak 설치 + realm", [
            "Keycloak release 다운로드 → /opt/iot-platform/keycloak",
            "Postgres datasource 설정 (KEYCLOAK_DATABASE_*)",
            "keycloak.service systemd unit",
            "iot-platform realm 생성",
            "7 role + test user 1명 생성"
        ]),
        ("Day 4", "Nginx + HTTPS", [
            "Nginx 설치 · iot-platform.conf 작성",
            "/auth → 8080 reverse proxy",
            "/api · / placeholder location",
            "certbot --nginx 으로 Let's Encrypt 발급",
            "systemd timer로 인증서 자동 갱신 확인"
        ]),
        ("Day 5", "환경 분리 + 절차서", [
            "/etc/iot-platform/{backend,worker,scheduler,mqtt,db,keycloak}.env 템플릿",
            "운영 절차서 (PHASE1_OPS.md) 작성 — 설치 명령 단계별",
            "장애 대응 절차 (서비스 재시작 · 로그 위치)"
        ]),
        ("Day 6", "smoke test 자동화", [
            "scripts/phase1_smoke.sh 작성",
            "  systemctl is-active 4개",
            "  curl /auth/.well-known/openid-configuration",
            "  psql SELECT 1",
            "  mosquitto_pub/sub 인증 테스트"
        ]),
        ("Day 7", "검증 + 문서 정리", [
            "smoke test 전체 실행 → PASS 확인",
            "체크리스트 마크다운 결과 commit",
            "Phase 2 인터뷰 준비 (남은 트랙: Gateway OS · scale · PoC 산업 등)"
        ]),
    ]
    cw = Inches(6.1); ch = Inches(2.0)
    for i, (d, t, items) in enumerate(days):
        col = i % 2
        row = i // 2
        # Day 7 is alone in last row — let it span
        if i == 6:
            x = Inches(0.5); y = Inches(1.1 + row * 2.05); w = Inches(12.3)
        else:
            x = Inches(0.5 + col * 6.27); y = Inches(1.1 + row * 2.05); w = cw
        add_round_rect(s, x, y, w, ch, SURFACE, line_color=LINE)
        add_rect(s, x, y, Inches(0.9), ch, NAVY)
        add_textbox(s, x + Inches(0.1), y + Inches(0.05), Inches(0.85), Inches(0.4),
                    d, size=12, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), y + Inches(0.5), Inches(0.85), Inches(1.4),
                    t, size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        add_multiline(s, x + Inches(1.05), y + Inches(0.1), w - Inches(1.15), ch - Inches(0.2),
                      [(it, {"size": 9.5}) for it in items], bullet=True, line_spacing=1.3)

# ---------- Slide 35: 라이선스 ----------
def slide_license():
    s = blank_slide(prs)
    add_header(s, "라이선스 검토", "권장 vs 주의 — 상용화 시 OSS Notice / SBOM 필수", 35, TOTAL)
    rec = [
        ("VerneMQ", "Apache 2.0", "권장"),
        ("PostgreSQL", "PostgreSQL License", "권장"),
        ("Keycloak", "Apache 2.0", "권장"),
        ("Apache ECharts", "Apache 2.0", "권장"),
        ("Nginx", "BSD-like", "권장"),
        ("FastAPI", "MIT", "권장"),
        ("React + Vite", "MIT", "권장"),
        ("Prometheus", "Apache 2.0", "선택"),
    ]
    avoid = [
        ("EMQX 최신", "BSL 계열 이슈 가능"),
        ("MinIO", "AGPLv3"),
        ("Grafana 고객용 노출", "AGPLv3"),
        ("Loki", "AGPL"),
        ("SWUpdate", "GPLv2"),
        ("TimescaleDB Community 일부", "Timescale License 혼재"),
        ("Docker Desktop", "조직 규모 구독 이슈"),
    ]
    add_round_rect(s, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(7.5), Inches(0.5), GREEN)
    add_textbox(s, Inches(0.7), Inches(1.13), Inches(7.0), Inches(0.42),
                "권장 컴포넌트", size=13, bold=True, color=TEXT_LIGHT)
    yy = Inches(1.75)
    for n, lic, judge in rec:
        add_round_rect(s, Inches(0.7), yy, Inches(7.1), Inches(0.55), SUBTLE, line_color=LINE)
        add_textbox(s, Inches(0.85), yy + Inches(0.05), Inches(2.5), Inches(0.45),
                    n, size=11, bold=True, color=NAVY)
        add_textbox(s, Inches(3.4), yy + Inches(0.05), Inches(3.0), Inches(0.45),
                    lic, size=10, color=TEXT_DARK, font=FONT_MONO)
        col = GREEN if judge == "권장" else MUTED
        add_textbox(s, Inches(6.5), yy + Inches(0.05), Inches(1.2), Inches(0.45),
                    judge, size=10, bold=True, color=col)
        yy += Inches(0.62)

    add_round_rect(s, Inches(8.2), Inches(1.1), Inches(4.65), Inches(5.95), SURFACE, line_color=LINE)
    add_rect(s, Inches(8.2), Inches(1.1), Inches(4.65), Inches(0.5), RED)
    add_textbox(s, Inches(8.4), Inches(1.13), Inches(4.3), Inches(0.42),
                "피하거나 주의", size=13, bold=True, color=TEXT_LIGHT)
    yy = Inches(1.75)
    for n, why in avoid:
        add_round_rect(s, Inches(8.4), yy, Inches(4.3), Inches(0.7), SUBTLE, line_color=LINE)
        add_textbox(s, Inches(8.55), yy + Inches(0.05), Inches(4.0), Inches(0.3),
                    n, size=11, bold=True, color=RED)
        add_textbox(s, Inches(8.55), yy + Inches(0.32), Inches(4.0), Inches(0.35),
                    why, size=9.5, color=TEXT_DARK)
        yy += Inches(0.78)

# ---------- Slide 36: 우선순위 ----------
def slide_priority():
    s = blank_slide(prs)
    add_header(s, "최종 구현 우선순위", "최우선 → 중간 → 제품화 — 자원 제한 시 위에서 아래로", 36, TOTAL)
    cols = [
        ("최우선",   AMBER, [
            "사용자별 다중 Gateway 권한 모델",
            "Gateway별 Sensor Profile / Channel",
            "MQTT topic + Gateway 인증",
            "Telemetry 저장 + latest 조회",
            "Gateway Config Versioning",
            "Command Request/Response",
        ]),
        ("중간",     NAVY, [
            "Sensor 추가 Wizard",
            "Dynamic Dashboard",
            "Alarm Rule",
            "Actuator Channel",
            "Config Rollback",
            "Audit Log",
        ]),
        ("제품화",   RED, [
            "Gateway별 X.509 인증서",
            "VerneMQ ACL",
            "OTA",
            "Backup / Restore",
            "Safety MCU 연동",
            "OSS Notice / SBOM",
        ]),
    ]
    cw = Inches(4.15); ch = Inches(5.95)
    for i, (n, c, items) in enumerate(cols):
        x = Inches(0.5 + i * 4.27)
        y = Inches(1.1)
        add_round_rect(s, x, y, cw, ch, SURFACE, line_color=LINE)
        add_rect(s, x, y, cw, Inches(0.65), c)
        add_textbox(s, x + Inches(0.2), y + Inches(0.1), cw, Inches(0.45),
                    n, size=15, bold=True, color=TEXT_LIGHT)
        add_multiline(s, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.5), ch - Inches(0.95),
                      [(t, {"size": 12, "space_after": 8}) for t in items], bullet=True,
                      line_spacing=1.5)

# ---------- Slide 37: 결론 ----------
def slide_conclusion():
    s = blank_slide(prs)
    add_header(s, "최종 결론", "Docker 없이도 자체 호스팅 IoT Gateway 플랫폼은 가능하다", 37, TOTAL)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
                "산업용 서버 운영에서는 systemd 기반 서비스 단위 관리가 오히려 안정적이다.",
                size=14, color=NAVY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.4),
                "단, 다음 7가지 조건은 반드시 설계에 반영해야 한다.",
                size=11, color=TEXT_DARK, italic=True)
    musts = [
        "한 사용자가 여러 IoT Gateway를 가질 수 있음",
        "Gateway마다 연결된 센서 종류가 다를 수 있음",
        "Gateway마다 릴레이·밸브·펌프 구성이 다를 수 있음",
        "관리회사는 전체 Gateway를 통합 관리해야 함",
        "일반 사용자는 본인에게 할당된 Gateway만 접근",
        "센서 종류 추가 시 코드 수정 최소화",
        "Gateway 설정은 서버에서 중앙 관리·버전 관리",
    ]
    for i, m in enumerate(musts):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(2.15 + row * 0.5)
        add_round_rect(s, x, y, Inches(0.45), Inches(0.4), AMBER)
        add_textbox(s, x, y + Inches(0.04), Inches(0.45), Inches(0.32),
                    str(i+1), size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.55), y + Inches(0.05), Inches(5.7), Inches(0.32),
                    m, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    add_round_rect(s, Inches(0.5), Inches(4.45), Inches(12.3), Inches(2.6), DARK, line_color=AMBER, line_w=1.5)
    add_textbox(s, Inches(0.7), Inches(4.55), Inches(12), Inches(0.4),
                "권장 최종 구조", size=13, bold=True, color=AMBER)
    final = ("VerneMQ  +  PostgreSQL  +  Keycloak  +  FastAPI Backend  +  React (Vite) Web Portal\n"
             "+  systemd 기반 서비스 운영  +  Gateway Profile  +  Sensor Profile  +  Sensor Channel Mapping\n"
             "+  Actuator Channel Mapping  +  Gateway Config Versioning  +  Dynamic Dashboard\n"
             "+  RBAC / ABAC 권한 모델  +  CM4 Linux + Safety MCU 이중 구조")
    add_textbox(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2.0), final,
                size=12.5, color=TEXT_LIGHT, anchor=MSO_ANCHOR.TOP)

# ---------- Slide 38: Next Steps ----------
def slide_next_steps():
    s = blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(0), Inches(4.0), Inches(7.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = DARK
    accent.line.fill.background(); accent.shadow.inherit = False
    add_rect(s, Inches(11.93), Inches(0), Inches(0.1), SLIDE_W, AMBER)

    add_rect(s, Inches(4.5), Inches(1.0), Inches(0.8), Inches(0.1), AMBER)
    add_textbox(s, Inches(4.5), Inches(1.2), Inches(8.5), Inches(0.7),
                "Next Steps", size=36, bold=True, color=TEXT_LIGHT)
    add_textbox(s, Inches(4.5), Inches(1.95), Inches(8.5), Inches(0.5),
                "지금 당장 시작 가능한 4가지", size=14, color=GOLD)

    nexts = [
        ("01", "Phase 1 sprint 시작",
         "본 자료의 Day 1-7 일정대로 서버 인프라 1주 구축. 코드 0줄."),
        ("02", "Phase 1 산출물 commit",
         "운영 절차서(PHASE1_OPS.md) + scripts/phase1_smoke.sh를 git repo로 형상 관리."),
        ("03", "Phase 2 인터뷰 준비",
         "남은 미결정: Gateway OS · v1 운영 규모 · 첫 PoC 산업 · Safety MCU 시점."),
        ("04", "Hardware lab 구성",
         "보유 CM4 Gateway + RS-485 온습도 센서 1개로 Phase 3 telemetry 흐름 실측."),
    ]
    y = Inches(2.7)
    for no, t, d in nexts:
        add_round_rect(s, Inches(4.5), y, Inches(8.3), Inches(0.95), DARK, line_color=AMBER)
        add_textbox(s, Inches(4.65), y + Inches(0.1), Inches(0.8), Inches(0.7),
                    no, size=22, bold=True, color=AMBER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(5.6), y + Inches(0.07), Inches(7.0), Inches(0.32),
                    t, size=14, bold=True, color=TEXT_LIGHT)
        add_textbox(s, Inches(5.6), y + Inches(0.42), Inches(7.0), Inches(0.5),
                    d, size=10.5, color=RGBColor(0xCC, 0xD4, 0xDB))
        y += Inches(1.05)

    add_textbox(s, Inches(4.5), Inches(7.0), Inches(8.5), Inches(0.3),
                "v1.0 · 작성: 2026-05-02 · Owner: Solo + Claude Code AI pair · CM4 IoT Gateway 자체 호스팅 플랫폼",
                size=9, color=MUTED, align=PP_ALIGN.LEFT)

# ===== Build all slides =====
slide_title()
slide_agenda()
slide_overview()
slide_premises()
slide_philosophy()
slide_concept_diagram()
slide_detailed_arch()
slide_systemd()
slide_directory()
slide_roles()
slide_rbac_hierarchy()
slide_keycloak()
slide_erd()
slide_core_tables()
slide_sensor_profile()
slide_telemetry()
slide_gw_config()
slide_mqtt_topic()
slide_mqtt_payload()
slide_command_flow()
slide_command_safety()
slide_web_portal()
slide_sensor_wizard()
slide_dynamic_dashboard()
slide_bulk_ops()
slide_gateway_agent()
slide_safety()
slide_api_design()
slide_security()
slide_rls()
slide_operations()
slide_roadmap()
slide_phase1_spec()
slide_phase1_schedule()
slide_license()
slide_priority()
slide_conclusion()
slide_next_steps()

OUTPUT = "/mnt/d/20_claude/IoT_Gateway_Server/CM4_IoT_Gateway_상세구현계획서.pptx"
prs.save(OUTPUT)
print(f"✅ PPTX 생성 완료: {OUTPUT}")
print(f"   슬라이드 수: {len(prs.slides)}")
